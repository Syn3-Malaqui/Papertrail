"""
Document Parser Module
Handles finding and extracting text from multiple document formats.
"""

import os
import glob
from pathlib import Path
from typing import List, Dict, Optional

# PDF parsing
from pdfminer.high_level import extract_text as pdf_extract_text
from pdfminer.pdfparser import PDFSyntaxError

# DOCX parsing
from docx import Document

# Additional format imports
try:
    import docx2txt  # For .doc files
except ImportError:
    docx2txt = None

try:
    from bs4 import BeautifulSoup  # For HTML files
except ImportError:
    BeautifulSoup = None

try:
    from striprtf.striprtf import rtf_to_text  # For RTF files
except ImportError:
    rtf_to_text = None

try:
    from odf import text, teletype
    from odf.opendocument import load as odf_load  # For ODT files
except ImportError:
    text = teletype = odf_load = None

try:
    from pptx import Presentation  # For PPTX files
except ImportError:
    Presentation = None

try:
    import openpyxl  # For XLSX files
    import xlrd      # For XLS files
    import pandas as pd
except ImportError:
    openpyxl = xlrd = pd = None

class DocumentParser:
    """Handles document discovery and text extraction."""
    
    SUPPORTED_EXTENSIONS = [
        '.pdf', '.txt', '.docx', '.doc',           # Original + DOC
        '.html', '.htm',                           # Web documents
        '.rtf',                                    # Rich Text Format
        '.odt',                                    # OpenDocument Text
        '.pptx', '.ppt',                          # PowerPoint
        '.xlsx', '.xls'                           # Excel
    ]
    
    def __init__(self):
        self.parsed_files = []
        self.errors = []
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check if optional dependencies are available and warn if missing."""
        missing_deps = []
        
        if docx2txt is None:
            missing_deps.append("docx2txt (for .doc files)")
        if BeautifulSoup is None:
            missing_deps.append("beautifulsoup4 (for .html files)")
        if rtf_to_text is None:
            missing_deps.append("striprtf (for .rtf files)")
        if text is None or odf_load is None:
            missing_deps.append("odfpy (for .odt files)")
        if Presentation is None:
            missing_deps.append("python-pptx (for .pptx files)")
        if openpyxl is None or xlrd is None:
            missing_deps.append("openpyxl/xlrd (for Excel files)")
        
        if missing_deps:
            print(f"⚠️  Warning: Some optional dependencies missing: {', '.join(missing_deps)}")
            print("   Install with: pip install -r requirements.txt")

    def find_documents(self, folder_path: str) -> List[str]:
        """
        Recursively find all supported document files in the given folder.
        
        Args:
            folder_path: Path to the folder to search
            
        Returns:
            List of file paths
        """
        if not os.path.exists(folder_path):
            raise FileNotFoundError(f"Folder not found: {folder_path}")
        
        file_paths = []
        
        for ext in self.SUPPORTED_EXTENSIONS:
            pattern = os.path.join(folder_path, '**', f'*{ext}')
            files = glob.glob(pattern, recursive=True)
            file_paths.extend(files)
        
        return sorted(file_paths)

    def extract_text_from_pdf(self, file_path: str) -> Optional[str]:
        """Extract text from PDF file."""
        try:
            text = pdf_extract_text(file_path)
            return text.strip() if text else None
        except (PDFSyntaxError, Exception) as e:
            self.errors.append(f"PDF parsing error for {file_path}: {str(e)}")
            return None

    def extract_text_from_docx(self, file_path: str) -> Optional[str]:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            paragraphs = [paragraph.text for paragraph in doc.paragraphs]
            text = '\n'.join(paragraphs)
            return text.strip() if text else None
        except Exception as e:
            self.errors.append(f"DOCX parsing error for {file_path}: {str(e)}")
            return None

    def extract_text_from_doc(self, file_path: str) -> Optional[str]:
        """Extract text from legacy DOC file."""
        if docx2txt is None:
            self.errors.append(f"DOC parsing requires docx2txt library: {file_path}")
            return None
        
        try:
            text = docx2txt.process(file_path)
            return text.strip() if text else None
        except Exception as e:
            self.errors.append(f"DOC parsing error for {file_path}: {str(e)}")
            return None

    def extract_text_from_html(self, file_path: str) -> Optional[str]:
        """Extract text from HTML file."""
        if BeautifulSoup is None:
            self.errors.append(f"HTML parsing requires beautifulsoup4 library: {file_path}")
            return None
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text and clean up whitespace
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return text.strip() if text else None
        except Exception as e:
            self.errors.append(f"HTML parsing error for {file_path}: {str(e)}")
            return None

    def extract_text_from_rtf(self, file_path: str) -> Optional[str]:
        """Extract text from RTF file."""
        if rtf_to_text is None:
            self.errors.append(f"RTF parsing requires striprtf library: {file_path}")
            return None
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                rtf_content = file.read()
            
            text = rtf_to_text(rtf_content)
            return text.strip() if text else None
        except Exception as e:
            self.errors.append(f"RTF parsing error for {file_path}: {str(e)}")
            return None

    def extract_text_from_odt(self, file_path: str) -> Optional[str]:
        """Extract text from OpenDocument Text file."""
        if text is None or odf_load is None:
            self.errors.append(f"ODT parsing requires odfpy library: {file_path}")
            return None
        
        try:
            doc = odf_load(file_path)
            all_paras = doc.getElementsByType(text.P)
            text_content = []
            
            for paragraph in all_paras:
                para_text = teletype.extractText(paragraph)
                if para_text.strip():
                    text_content.append(para_text.strip())
            
            result = '\n'.join(text_content)
            return result.strip() if result else None
        except Exception as e:
            self.errors.append(f"ODT parsing error for {file_path}: {str(e)}")
            return None

    def extract_text_from_pptx(self, file_path: str) -> Optional[str]:
        """Extract text from PowerPoint PPTX file."""
        if Presentation is None:
            self.errors.append(f"PPTX parsing requires python-pptx library: {file_path}")
            return None
        
        try:
            presentation = Presentation(file_path)
            text_content = []
            
            for slide in presentation.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_content.append(' '.join(slide_text))
            
            result = '\n'.join(text_content)
            return result.strip() if result else None
        except Exception as e:
            self.errors.append(f"PPTX parsing error for {file_path}: {str(e)}")
            return None

    def extract_text_from_ppt(self, file_path: str) -> Optional[str]:
        """Extract text from legacy PowerPoint PPT file."""
        # Note: Legacy PPT support is limited; consider using LibreOffice conversion
        self.errors.append(f"Legacy PPT format not fully supported: {file_path}")
        return None

    def extract_text_from_xlsx(self, file_path: str) -> Optional[str]:
        """Extract text from Excel XLSX file."""
        if openpyxl is None or pd is None:
            self.errors.append(f"XLSX parsing requires openpyxl and pandas libraries: {file_path}")
            return None
        
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_content = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert all data to strings and join
                sheet_text = []
                for column in df.columns:
                    col_data = df[column].dropna().astype(str).tolist()
                    sheet_text.extend(col_data)
                
                if sheet_text:
                    text_content.append(f"Sheet {sheet_name}: " + ' '.join(sheet_text))
            
            result = '\n'.join(text_content)
            return result.strip() if result else None
        except Exception as e:
            self.errors.append(f"XLSX parsing error for {file_path}: {str(e)}")
            return None

    def extract_text_from_xls(self, file_path: str) -> Optional[str]:
        """Extract text from legacy Excel XLS file."""
        if xlrd is None or pd is None:
            self.errors.append(f"XLS parsing requires xlrd and pandas libraries: {file_path}")
            return None
        
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path, engine='xlrd')
            text_content = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name, engine='xlrd')
                
                # Convert all data to strings and join
                sheet_text = []
                for column in df.columns:
                    col_data = df[column].dropna().astype(str).tolist()
                    sheet_text.extend(col_data)
                
                if sheet_text:
                    text_content.append(f"Sheet {sheet_name}: " + ' '.join(sheet_text))
            
            result = '\n'.join(text_content)
            return result.strip() if result else None
        except Exception as e:
            self.errors.append(f"XLS parsing error for {file_path}: {str(e)}")
            return None

    def extract_text_from_txt(self, file_path: str) -> Optional[str]:
        """Extract text from TXT file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
            return text.strip() if text else None
        except Exception as e:
            self.errors.append(f"TXT parsing error for {file_path}: {str(e)}")
            return None

    def extract_text(self, file_path: str) -> Optional[str]:
        """
        Extract text from a file based on its extension.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Extracted text or None if failed
        """
        file_ext = Path(file_path).suffix.lower()
        
        # Route to appropriate extraction method
        if file_ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif file_ext == '.docx':
            return self.extract_text_from_docx(file_path)
        elif file_ext == '.doc':
            return self.extract_text_from_doc(file_path)
        elif file_ext in ['.html', '.htm']:
            return self.extract_text_from_html(file_path)
        elif file_ext == '.rtf':
            return self.extract_text_from_rtf(file_path)
        elif file_ext == '.odt':
            return self.extract_text_from_odt(file_path)
        elif file_ext == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif file_ext == '.ppt':
            return self.extract_text_from_ppt(file_path)
        elif file_ext == '.xlsx':
            return self.extract_text_from_xlsx(file_path)
        elif file_ext == '.xls':
            return self.extract_text_from_xls(file_path)
        elif file_ext == '.txt':
            return self.extract_text_from_txt(file_path)
        else:
            self.errors.append(f"Unsupported file type: {file_ext}")
            return None

    def parse_documents(self, folder_path: str) -> Dict[str, str]:
        """
        Parse all documents in a folder and return a mapping of filenames to text.
        
        Args:
            folder_path: Path to the folder containing documents
            
        Returns:
            Dictionary mapping filenames to extracted text
        """
        self.parsed_files = []
        self.errors = []
        
        file_paths = self.find_documents(folder_path)
        results = {}
        
        print(f"Found {len(file_paths)} documents to process...")
        print(f"Supported formats: {', '.join(self.SUPPORTED_EXTENSIONS)}")
        
        for file_path in file_paths:
            print(f"Processing: {os.path.basename(file_path)}")
            
            text = self.extract_text(file_path)
            if text:
                filename = os.path.basename(file_path)
                results[filename] = text
                self.parsed_files.append(file_path)
            else:
                print(f"Failed to extract text from: {file_path}")
        
        if self.errors:
            print(f"\nEncountered {len(self.errors)} errors during parsing:")
            for error in self.errors:
                print(f"  - {error}")
        
        print(f"\nSuccessfully parsed {len(results)} documents.")
        return results

# Example usage
if __name__ == "__main__":
    parser = DocumentParser()
    
    # Test with a sample directory
    test_folder = input("Enter folder path to test (or press Enter for current directory): ").strip()
    if not test_folder:
        test_folder = "."
    
    try:
        documents = parser.parse_documents(test_folder)
        
        if documents:
            print("\n--- Sample Results ---")
            for filename, text in list(documents.items())[:2]:  # Show first 2 files
                print(f"\nFile: {filename}")
                print(f"Text preview: {text[:200]}...")
        else:
            print("No documents found or successfully parsed.")
            
    except Exception as e:
        print(f"Error: {e}") 