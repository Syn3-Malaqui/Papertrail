#!/usr/bin/env python3
"""
Diverse Sample Document Generator for Papertrail
Creates realistic documents in multiple formats with proper formatting.
Now supports: PDF, TXT, DOCX, DOC, HTML/HTM, RTF, ODT, PPTX, PPT, XLSX, XLS
"""

import os
import random
import io
from datetime import datetime, timedelta
from typing import List, Dict, Tuple

# PDF generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors

# DOCX generation
from docx import Document as DocxDocument
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Additional format imports (with graceful fallbacks)
try:
    import pandas as pd
    import openpyxl
    from openpyxl.styles import Font, Alignment
except ImportError:
    pd = openpyxl = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
except ImportError:
    Presentation = PptxInches = None

try:
    from odf.opendocument import OpenDocumentText
    from odf.style import Style, TextProperties, ParagraphProperties
    from odf.text import P, H
except ImportError:
    OpenDocumentText = Style = TextProperties = ParagraphProperties = P = H = None

class DiverseDocumentGenerator:
    """Generates realistic documents in multiple formats with proper formatting."""
    
    def __init__(self, output_dir: str = "diverse_sample_documents"):
        self.output_dir = output_dir
        self.setup_data()
        self._check_dependencies()
        
    def _check_dependencies(self):
        """Check and warn about missing dependencies."""
        missing = []
        if pd is None or openpyxl is None:
            missing.append("pandas/openpyxl (for Excel files)")
        if Presentation is None:
            missing.append("python-pptx (for PowerPoint files)")
        if OpenDocumentText is None:
            missing.append("odfpy (for ODT files)")
        
        if missing:
            print(f"⚠️  Warning: Missing dependencies for: {', '.join(missing)}")
            print("   Install with: pip install -r requirements.txt")
        
    def setup_data(self):
        """Setup realistic data for generating documents."""
        self.companies = [
            "TechnoGlobal Solutions Inc", "Meridian Consulting Group", "Apex Digital Systems",
            "Pinnacle Financial Services", "Innovative Research Labs", "Strategic Partners LLC",
            "NextGen Technologies Corp", "Elite Business Solutions", "ProActive Enterprises",
            "Dynamic Consulting Group", "Premier Analytics Inc", "Advanced Systems Ltd",
            "Innovation Dynamics Corp", "Strategic Alliance Group", "Digital Transform LLC",
            "Enterprise Solutions Inc", "Global Systems Technologies", "Future Vision Consulting"
        ]
        
        self.people = [
            ("John", "Mitchell", "CEO"), ("Sarah", "Chen", "CFO"), ("Michael", "Rodriguez", "CTO"),
            ("Emily", "Thompson", "VP Operations"), ("David", "Kumar", "Legal Counsel"),
            ("Jessica", "Anderson", "HR Director"), ("Robert", "Williams", "Sales Manager"),
            ("Amanda", "Taylor", "Project Manager"), ("James", "Brown", "Senior Analyst"),
            ("Lauren", "Davis", "Marketing Director"), ("Christopher", "Wilson", "Finance Manager"),
            ("Michelle", "Garcia", "Operations Lead"), ("Daniel", "Lee", "Technical Lead"),
            ("Rachel", "Martinez", "Business Analyst"), ("Kevin", "Johnson", "Quality Manager"),
            ("Lisa", "Zhang", "Data Scientist"), ("Mark", "Thompson", "Security Officer")
        ]
        
        self.addresses = [
            ("1250 Tech Park Drive", "San Jose", "CA", "95110"),
            ("890 Business Center Blvd", "Austin", "TX", "78701"),
            ("2100 Corporate Square", "New York", "NY", "10001"),
            ("555 Innovation Way", "Seattle", "WA", "98101"),
            ("1775 Professional Plaza", "Chicago", "IL", "60601"),
            ("3300 Executive Circle", "Dallas", "TX", "75201"),
            ("4400 Technology Center", "Boston", "MA", "02101"),
            ("5500 Corporate Drive", "Atlanta", "GA", "30301")
        ]

        # Enhanced document content templates
        self.invoice_services = [
            "Professional Consulting Services", "Software Development", "System Integration",
            "Technical Support Services", "Project Management", "Data Analysis Services",
            "Training and Education", "Quality Assurance Testing", "Security Assessment"
        ]
        
        self.memo_subjects = [
            "Policy Update - Remote Work Guidelines", "Quarterly Team Meeting Schedule",
            "New Safety Procedures Implementation", "Budget Allocation for Q4 Projects",
            "Performance Review Process Updates", "Holiday Schedule Announcement",
            "Training Program Requirements", "Office Relocation Timeline",
            "IT Security Protocol Changes", "Employee Benefits Enhancement"
        ]
        
        self.contract_types = [
            "Professional Services Agreement", "Software License Agreement", 
            "Consulting Services Contract", "Maintenance and Support Agreement",
            "Employment Agreement", "Non-Disclosure Agreement",
            "Service Level Agreement", "Partnership Agreement"
        ]

    def create_output_dir(self):
        """Create output directory structure for all supported formats."""
        os.makedirs(self.output_dir, exist_ok=True)
        formats = ['pdf', 'docx', 'txt', 'doc', 'html', 'rtf', 'odt', 'pptx', 'ppt', 'xlsx', 'xls']
        for fmt in formats:
            os.makedirs(os.path.join(self.output_dir, fmt), exist_ok=True)

    def random_date(self, days_back: int = 365) -> datetime:
        """Generate random date."""
        return datetime.now() - timedelta(days=random.randint(0, days_back))

    def format_currency(self, amount: float) -> str:
        """Format currency amount."""
        return f"${amount:,.2f}"

    # ===========================================
    # NEW FILE FORMAT GENERATORS
    # ===========================================

    def generate_html_document(self, filename: str, doc_type: str, content_data: dict) -> str:
        """Generate HTML document."""
        filepath = os.path.join(self.output_dir, 'html', filename)
        
        html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{content_data['title']}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
        .header {{ text-align: center; border-bottom: 2px solid #333; padding-bottom: 20px; }}
        .company {{ font-size: 24px; font-weight: bold; color: #2c3e50; }}
        .doc-type {{ font-size: 18px; margin: 10px 0; color: #34495e; }}
        .content {{ margin: 30px 0; }}
        .section {{ margin: 20px 0; }}
        .section h3 {{ color: #2980b9; border-bottom: 1px solid #3498db; }}
        table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
        th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
        th {{ background-color: #f2f2f2; }}
        .footer {{ margin-top: 40px; border-top: 1px solid #ddd; padding-top: 20px; }}
    </style>
</head>
<body>
    <div class="header">
        <div class="company">{content_data['company']}</div>
        <div class="doc-type">{content_data['title']}</div>
        <div>Date: {content_data['date']}</div>
    </div>
    
    <div class="content">
        {content_data['body']}
    </div>
    
    <div class="footer">
        <p><strong>Document Information:</strong></p>
        <p>Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</p>
        <p>Author: {content_data.get('author', 'System Generated')}</p>
        <p>Version: {random.randint(1, 5)}.{random.randint(0, 9)}</p>
    </div>
</body>
</html>"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(html_content)
        return filepath

    def generate_rtf_document(self, filename: str, doc_type: str, content_data: dict) -> str:
        """Generate RTF document."""
        filepath = os.path.join(self.output_dir, 'rtf', filename)
        
        # Basic RTF structure
        rtf_content = r"""{{\rtf1\ansi\ansicpg1252\deff0 {{\fonttbl{{\f0\fswiss\fcharset0 Arial;}}}}
{{\colortbl;\red0\green0\blue0;\red46\green125\blue186;}}
\viewkind4\uc1\pard\sa200\sl276\slmult1\qc\cf2\b\f0\fs28 """ + content_data['company'] + r"""
\par """ + content_data['title'] + r"""
\par \cf1\b0\fs20 Date: """ + content_data['date'] + r"""
\par \pard\sa200\sl276\slmult1\ql
\par """ + content_data['body'].replace('\n', r'\par ') + r"""
\par 
\par \b Document Information:\b0
\par Generated: """ + datetime.now().strftime('%B %d, %Y') + r"""
\par Author: """ + content_data.get('author', 'System Generated') + r"""
\par Version: """ + f"{random.randint(1, 5)}.{random.randint(0, 9)}" + r"""
\par }}"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(rtf_content)
        return filepath

    def generate_odt_document(self, filename: str, doc_type: str, content_data: dict) -> str:
        """Generate ODT document."""
        if OpenDocumentText is None:
            return None
            
        filepath = os.path.join(self.output_dir, 'odt', filename)
        
        try:
            doc = OpenDocumentText()
            
            # Add title
            title = H(outlinelevel=1, text=content_data['title'])
            doc.text.addElement(title)
            
            # Add company and date
            company_para = P(text=f"Company: {content_data['company']}")
            doc.text.addElement(company_para)
            
            date_para = P(text=f"Date: {content_data['date']}")
            doc.text.addElement(date_para)
            
            # Add empty line
            doc.text.addElement(P())
            
            # Add content (split by lines)
            for line in content_data['body'].split('\n'):
                if line.strip():
                    para = P(text=line.strip())
                    doc.text.addElement(para)
                else:
                    doc.text.addElement(P())
            
            # Add footer info
            doc.text.addElement(P())
            footer_info = P(text=f"Generated: {datetime.now().strftime('%B %d, %Y')} | Author: {content_data.get('author', 'System Generated')}")
            doc.text.addElement(footer_info)
            
            doc.save(filepath)
            return filepath
        except Exception as e:
            print(f"Error generating ODT: {e}")
            return None

    def generate_pptx_document(self, filename: str, doc_type: str, content_data: dict) -> str:
        """Generate PowerPoint PPTX document."""
        if Presentation is None:
            return None
            
        filepath = os.path.join(self.output_dir, 'pptx', filename)
        
        try:
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = content_data['title']
            subtitle.text = f"{content_data['company']}\n{content_data['date']}"
            
            # Content slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Overview"
            
            # Split content into bullet points
            content_lines = content_data['body'].split('\n')
            bullet_text = ""
            for line in content_lines[:8]:  # Limit to 8 bullet points
                if line.strip():
                    bullet_text += f"• {line.strip()}\n"
            
            content.text = bullet_text
            
            # Summary slide if content is long
            if len(content_lines) > 8:
                summary_slide = prs.slides.add_slide(bullet_slide_layout)
                summary_slide.shapes.title.text = "Additional Information"
                summary_content = summary_slide.placeholders[1]
                
                remaining_text = ""
                for line in content_lines[8:15]:  # Next 7 lines
                    if line.strip():
                        remaining_text += f"• {line.strip()}\n"
                summary_content.text = remaining_text
            
            prs.save(filepath)
            return filepath
        except Exception as e:
            print(f"Error generating PPTX: {e}")
            return None

    def generate_xlsx_document(self, filename: str, doc_type: str, content_data: dict) -> str:
        """Generate Excel XLSX document."""
        if pd is None or openpyxl is None:
            return None
            
        filepath = os.path.join(self.output_dir, 'xlsx', filename)
        
        try:
            with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
                # Create different sheets based on document type
                if 'invoice' in doc_type.lower():
                    # Invoice data
                    invoice_data = {
                        'Description': [
                            'Professional Consulting Services',
                            'Technical Support',
                            'Project Management',
                            'Quality Assurance',
                            'Documentation',
                            '',
                            'Subtotal',
                            'Tax (8.75%)',
                            'Total Amount Due'
                        ],
                        'Hours': [40, 20, 15, 10, 5, '', '', '', ''],
                        'Rate': [150, 125, 175, 100, 75, '', '', '', ''],
                        'Amount': [6000, 2500, 2625, 1000, 375, '', 12500, 1093.75, 13593.75]
                    }
                    df = pd.DataFrame(invoice_data)
                    df.to_excel(writer, sheet_name='Invoice Details', index=False)
                    
                elif 'report' in doc_type.lower():
                    # Report data with multiple sheets
                    # Financial data
                    financial_data = {
                        'Quarter': ['Q1 2024', 'Q2 2024', 'Q3 2024', 'Q4 2024'],
                        'Revenue': [450000, 520000, 480000, 610000],
                        'Expenses': [320000, 340000, 350000, 380000],
                        'Profit': [130000, 180000, 130000, 230000],
                        'Growth %': [5.2, 15.6, -7.7, 27.1]
                    }
                    df_financial = pd.DataFrame(financial_data)
                    df_financial.to_excel(writer, sheet_name='Financial Summary', index=False)
                    
                    # Performance metrics
                    metrics_data = {
                        'Metric': ['Customer Satisfaction', 'Employee Retention', 'Market Share', 
                                 'Response Time', 'Quality Score', 'Cost Efficiency'],
                        'Current': [4.2, 92, 15.8, 2.3, 94.5, 87.2],
                        'Target': [4.5, 95, 18.0, 2.0, 95.0, 90.0],
                        'Status': ['Below Target', 'Below Target', 'Below Target', 
                                 'Above Target', 'Below Target', 'Below Target']
                    }
                    df_metrics = pd.DataFrame(metrics_data)
                    df_metrics.to_excel(writer, sheet_name='Performance Metrics', index=False)
                    
                else:
                    # General document data
                    doc_info = {
                        'Field': ['Document Type', 'Company', 'Date', 'Author', 'Version', 'Status'],
                        'Value': [content_data['title'], content_data['company'], 
                                content_data['date'], content_data.get('author', 'System Generated'),
                                f"{random.randint(1, 5)}.{random.randint(0, 9)}", 'Active']
                    }
                    df_info = pd.DataFrame(doc_info)
                    df_info.to_excel(writer, sheet_name='Document Info', index=False)
                    
                    # Sample data table
                    sample_data = {
                        'Item': [f'Item {i+1}' for i in range(10)],
                        'Category': [random.choice(['A', 'B', 'C']) for _ in range(10)],
                        'Value': [random.randint(100, 1000) for _ in range(10)],
                        'Status': [random.choice(['Active', 'Inactive', 'Pending']) for _ in range(10)]
                    }
                    df_sample = pd.DataFrame(sample_data)
                    df_sample.to_excel(writer, sheet_name='Sample Data', index=False)
            
            return filepath
        except Exception as e:
            print(f"Error generating XLSX: {e}")
            return None

    def generate_xls_document(self, filename: str, doc_type: str, content_data: dict) -> str:
        """Generate legacy Excel XLS document (simplified version)."""
        if pd is None:
            return None
            
        filepath = os.path.join(self.output_dir, 'xls', filename)
        
        try:
            # Create simple data for XLS
            if 'invoice' in doc_type.lower():
                data = {
                    'Item': ['Consulting Services', 'Technical Support', 'Training'],
                    'Amount': [5000.00, 2500.00, 1500.00],
                    'Tax': [437.50, 218.75, 131.25],
                    'Total': [5437.50, 2718.75, 1631.25]
                }
            elif 'report' in doc_type.lower():
                data = {
                    'Month': ['January', 'February', 'March', 'April'],
                    'Revenue': [45000, 52000, 48000, 61000],
                    'Expenses': [32000, 34000, 35000, 38000],
                    'Profit': [13000, 18000, 13000, 23000]
                }
            else:
                data = {
                    'Field': ['Document', 'Company', 'Date', 'Version'],
                    'Value': [content_data['title'], content_data['company'], 
                            content_data['date'], f"{random.randint(1, 5)}.{random.randint(0, 9)}"]
                }
            
            df = pd.DataFrame(data)
            df.to_excel(filepath, index=False, engine='xlwt')
            return filepath
        except Exception as e:
            print(f"Error generating XLS: {e}")
            return None

    def generate_doc_document(self, filename: str, doc_type: str, content_data: dict) -> str:
        """Generate legacy DOC document (creates RTF with .doc extension for compatibility)."""
        # Since true .doc generation is complex, we'll create RTF content with .doc extension
        # This will be readable by most Word processors
        filepath = os.path.join(self.output_dir, 'doc', filename)
        
        rtf_content = r"""{{\rtf1\ansi\ansicpg1252\deff0 {{\fonttbl{{\f0\fswiss\fcharset0 Arial;}}}}
{{\colortbl;\red0\green0\blue0;\red46\green125\blue186;}}
\viewkind4\uc1\pard\sa200\sl276\slmult1\qc\cf2\b\f0\fs28 """ + content_data['company'] + r"""
\par """ + content_data['title'] + r"""
\par \cf1\b0\fs20 Date: """ + content_data['date'] + r"""
\par \pard\sa200\sl276\slmult1\ql
\par """ + content_data['body'].replace('\n', r'\par ') + r"""
\par 
\par \b Document Information:\b0
\par Generated: """ + datetime.now().strftime('%B %d, %Y') + r"""
\par Author: """ + content_data.get('author', 'System Generated') + r"""
\par Format: Legacy DOC (RTF Compatible)
\par }}"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(rtf_content)
        return filepath

    # ===========================================
    # ENHANCED CONTENT GENERATORS FOR EACH CATEGORY
    # ===========================================

    def create_invoice_content(self, fmt: str) -> dict:
        """Create enhanced invoice content."""
        company = random.choice(self.companies)
        invoice_num = f"INV-{random.randint(2023, 2024)}-{random.randint(1000, 9999)}"
        date = self.random_date(90).strftime('%B %d, %Y')
        client_company = random.choice([c for c in self.companies if c != company])
        
        services = random.sample(self.invoice_services, random.randint(2, 4))
        subtotal = random.randint(5000, 25000)
        tax_rate = 0.0875
        tax = subtotal * tax_rate
        total = subtotal + tax
        
        if fmt in ['html', 'rtf', 'odt']:
            body = f"""
            <div class="section">
                <h3>Invoice Details</h3>
                <p><strong>Invoice Number:</strong> {invoice_num}</p>
                <p><strong>Bill To:</strong> {client_company}</p>
                <p><strong>Due Date:</strong> {(datetime.now() + timedelta(days=30)).strftime('%B %d, %Y')}</p>
            </div>
            
            <div class="section">
                <h3>Services Provided</h3>
                <table>
                    <tr><th>Description</th><th>Amount</th></tr>
                    {"".join([f"<tr><td>{service}</td><td>${random.randint(1000, 8000):,}.00</td></tr>" for service in services])}
                    <tr><th>Subtotal</th><th>${subtotal:,}.00</th></tr>
                    <tr><th>Tax (8.75%)</th><th>${tax:,.2f}</th></tr>
                    <tr><th>Total Amount Due</th><th>${total:,.2f}</th></tr>
                </table>
            </div>
            
            <div class="section">
                <h3>Payment Information</h3>
                <p>Payment Terms: Net 30 days</p>
                <p>Please remit payment to the address above or contact our billing department.</p>
            </div>
            """
        else:
            body = f"""Invoice Number: {invoice_num}
Bill To: {client_company}
Due Date: {(datetime.now() + timedelta(days=30)).strftime('%B %d, %Y')}

SERVICES PROVIDED:
{chr(10).join([f'• {service}: ${random.randint(1000, 8000):,}.00' for service in services])}

Subtotal: ${subtotal:,}.00
Tax (8.75%): ${tax:,.2f}
TOTAL AMOUNT DUE: ${total:,.2f}

Payment Terms: Net 30 days
Please remit payment to the address above or contact our billing department for questions."""
        
        return {
            'title': 'PROFESSIONAL SERVICES INVOICE',
            'company': company,
            'date': date,
            'body': body,
            'author': f"Billing Department"
        }

    def generate_invoice_pdf(self, filename: str) -> str:
        """Generate a professional invoice PDF."""
        filepath = os.path.join(self.output_dir, 'pdf', filename)
        doc = SimpleDocTemplate(filepath, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            textColor=colors.darkblue
        )
        
        company = random.choice(self.companies)
        invoice_num = f"INV-{random.randint(2023, 2024)}-{random.randint(1000, 9999)}"
        date = self.random_date(90)
        client_company = random.choice([c for c in self.companies if c != company])
        
        # Header
        story.append(Paragraph(f"<b>{company}</b>", title_style))
        story.append(Paragraph("Professional Services Invoice", styles['Heading2']))
        story.append(Spacer(1, 20))
        
        # Invoice details
        details = [
            ['Invoice Number:', invoice_num],
            ['Date:', date.strftime('%B %d, %Y')],
            ['Due Date:', (date + timedelta(days=30)).strftime('%B %d, %Y')],
            ['Bill To:', client_company]
        ]
        
        details_table = Table(details, colWidths=[2*inch, 3*inch])
        details_table.setStyle(TableStyle([
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ]))
        story.append(details_table)
        story.append(Spacer(1, 30))
        
        # Services table
        services = [
            ['Description', 'Quantity', 'Rate', 'Amount'],
            ['Professional Consulting Services', '40 hrs', '$150.00', '$6,000.00'],
            ['Technical Documentation', '10 hrs', '$120.00', '$1,200.00'],
            ['Project Management', '20 hrs', '$130.00', '$2,600.00']
        ]
        
        subtotal = 9800.00
        tax = subtotal * 0.0875
        total = subtotal + tax
        
        services.extend([
            ['', '', 'Subtotal:', f'${subtotal:,.2f}'],
            ['', '', 'Sales Tax (8.75%):', f'${tax:,.2f}'],
            ['', '', 'TOTAL DUE:', f'${total:,.2f}']
        ])
        
        services_table = Table(services, colWidths=[3*inch, 1*inch, 1*inch, 1.5*inch])
        services_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
            ('BACKGROUND', (0, 1), (-1, 3), colors.beige),
            ('GRID', (0, 0), (-1, 3), 1, colors.black),
            ('FONTNAME', (2, -3), (-1, -1), 'Helvetica-Bold'),
            ('BACKGROUND', (2, -1), (-1, -1), colors.lightgrey)
        ]))
        
        story.append(services_table)
        story.append(Spacer(1, 30))
        
        # Payment terms
        story.append(Paragraph("<b>Payment Terms:</b> Net 30 days", styles['Normal']))
        story.append(Paragraph("Please remit payment within 30 days of invoice date.", styles['Normal']))
        story.append(Paragraph(f"Account Number: {random.randint(100000, 999999)}", styles['Normal']))
        
        doc.build(story)
        return filepath

    def generate_memo_docx(self, filename: str) -> str:
        """Generate a professional memo DOCX."""
        filepath = os.path.join(self.output_dir, 'docx', filename)
        doc = DocxDocument()
        
        # Header
        header = doc.sections[0].header
        header_para = header.paragraphs[0]
        header_para.text = random.choice(self.companies)
        header_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Title
        title = doc.add_heading('INTERNAL MEMORANDUM', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Memo header info
        sender_name, sender_last, sender_title = random.choice(self.people)
        memo_date = self.random_date(30)
        
        memo_subjects = [
            "Updated Remote Work Policy Implementation",
            "Q4 Budget Planning and Resource Allocation", 
            "New Employee Onboarding Process Changes",
            "Annual Performance Review Procedures",
            "IT Security Protocol Updates",
            "Office Space Renovation Schedule"
        ]
        
        subject = random.choice(memo_subjects)
        
        # Add memo details
        doc.add_paragraph(f"TO: All Department Managers")
        doc.add_paragraph(f"FROM: {sender_name} {sender_last}, {sender_title}")
        doc.add_paragraph(f"DATE: {memo_date.strftime('%B %d, %Y')}")
        doc.add_paragraph(f"RE: {subject}")
        doc.add_paragraph("")  # Blank line
        
        # Body content
        body_paragraphs = [
            f"This memorandum serves to inform all department managers about important updates regarding {subject.lower()}.",
            "",
            "Key Points:",
            f"• Effective Date: {(memo_date + timedelta(days=14)).strftime('%B %d, %Y')}",
            f"• Implementation Timeline: {random.choice(['2 weeks', '30 days', 'immediate'])}",
            f"• Required Actions: {random.choice(['Complete mandatory training', 'Submit compliance forms', 'Attend briefing sessions'])}",
            f"• Compliance Deadline: {(memo_date + timedelta(days=45)).strftime('%B %d, %Y')}",
            "",
            "Additional Information:",
            random.choice([
                "All affected employees must complete the required training modules by the specified deadline. Failure to comply may result in disciplinary action.",
                "Please review the attached documentation thoroughly and ensure your team follows the new procedures outlined in this memo.",
                "Department heads are responsible for cascading this information to their respective teams and monitoring compliance.",
                "For questions or clarification regarding these changes, please contact the HR department at extension 5500."
            ]),
            "",
            f"Thank you for your attention to this matter. For additional questions, please contact {sender_title} at extension {random.randint(3000, 5999)}.",
            "",
            "Best regards,",
            f"{sender_name} {sender_last}",
            f"{sender_title}"
        ]
        
        for para_text in body_paragraphs:
            para = doc.add_paragraph(para_text)
            if para_text.startswith("•"):
                para.style = 'List Bullet'
        
        doc.save(filepath)
        return filepath

    def generate_contract_txt(self, filename: str) -> str:
        """Generate a professional contract TXT with proper formatting."""
        filepath = os.path.join(self.output_dir, 'txt', filename)
        
        contract_types = [
            "Professional Services Agreement", "Software Development Contract",
            "Consulting Services Agreement", "Non-Disclosure Agreement",
            "Employment Contract", "Vendor Services Agreement"
        ]
        
        contract_type = random.choice(contract_types)
        provider = random.choice(self.companies)
        client = random.choice([c for c in self.companies if c != provider])
        date = self.random_date(60)
        provider_addr = random.choice(self.addresses)
        client_addr = random.choice([a for a in self.addresses if a != provider_addr])
        
        content = f"""
{'='*80}
{contract_type.upper()}
{'='*80}

This {contract_type} ("Agreement") is entered into on {date.strftime('%B %d, %Y')}, 
between {provider} ("Provider") and {client} ("Client").

PROVIDER INFORMATION:
{provider}
{provider_addr[0]}
{provider_addr[1]}, {provider_addr[2]} {provider_addr[3]}

CLIENT INFORMATION:
{client}
{client_addr[0]}
{client_addr[1]}, {client_addr[2]} {client_addr[3]}

{'='*80}
TERMS AND CONDITIONS
{'='*80}

1. SCOPE OF SERVICES
   Provider agrees to provide professional {random.choice(['consulting', 'development', 'analytical', 'technical'])} 
   services as outlined in Schedule A, attached hereto and incorporated by reference.
   
   Services include but are not limited to:
   • {random.choice(['Strategic planning and analysis', 'Software development and maintenance', 'Technical consulting and support'])}
   • {random.choice(['Project management and coordination', 'Quality assurance and testing', 'Documentation and training'])}
   • {random.choice(['Risk assessment and mitigation', 'Performance optimization', 'Compliance and regulatory support'])}

2. TERM AND TERMINATION
   This Agreement shall commence on {date.strftime('%B %d, %Y')} and shall continue 
   for a period of {random.choice([12, 18, 24, 36])} months, unless terminated earlier 
   in accordance with the provisions herein.
   
   Either party may terminate this Agreement with {random.choice([30, 60, 90])} days 
   written notice to the other party.

3. COMPENSATION
   Client agrees to pay Provider a total fee of {self.format_currency(random.randint(50000, 500000))} 
   for the services rendered under this Agreement.
   
   Payment Terms:
   • {random.choice(['Monthly', 'Quarterly', 'Milestone-based'])} payments
   • Net {random.choice([15, 30, 45])} days from invoice date
   • Late payment penalty: 1.5% per month

4. INTELLECTUAL PROPERTY
   All work products, deliverables, and intellectual property created under this 
   Agreement shall be owned by {random.choice(['Client', 'Provider', 'jointly by both parties'])}.

5. CONFIDENTIALITY
   Both parties acknowledge that they may have access to confidential information. 
   Each party agrees to maintain the confidentiality of such information and not 
   to disclose it to third parties without prior written consent.

6. WARRANTIES AND REPRESENTATIONS
   Provider represents and warrants that:
   • It has the authority to enter into this Agreement
   • The services will be performed in a professional manner
   • All work will comply with applicable laws and regulations

7. LIMITATION OF LIABILITY
   In no event shall either party be liable for any indirect, incidental, special, 
   or consequential damages arising out of this Agreement.

8. GOVERNING LAW
   This Agreement shall be governed by and construed in accordance with the laws 
   of the State of {random.choice(['California', 'New York', 'Texas', 'Delaware'])}.

9. ENTIRE AGREEMENT
   This Agreement constitutes the entire agreement between the parties and supersedes 
   all prior negotiations, representations, or agreements relating to the subject matter.

{'='*80}
SIGNATURES
{'='*80}

IN WITNESS WHEREOF, the parties have executed this Agreement as of the date first 
written above.

PROVIDER:                           CLIENT:
{provider:<35} {client}

By: _________________________      By: _________________________
Name: {random.choice(self.people)[0]} {random.choice(self.people)[1]:<22} Name: {random.choice(self.people)[0]} {random.choice(self.people)[1]}
Title: {random.choice(['CEO', 'President', 'COO']):<24} Title: {random.choice(['CEO', 'President', 'COO'])}
Date: ______________________      Date: ______________________

{'='*80}
END OF AGREEMENT
{'='*80}
"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return filepath

    def generate_legal_pdf(self, filename: str) -> str:
        """Generate a legal document PDF."""
        filepath = os.path.join(self.output_dir, 'pdf', filename)
        doc = SimpleDocTemplate(filepath, pagesize=A4)
        story = []
        styles = getSampleStyleSheet()
        
        # Legal header style
        legal_style = ParagraphStyle(
            'LegalHeader',
            parent=styles['Heading1'],
            fontSize=16,
            spaceAfter=20,
            alignment=1,  # Center
            textColor=colors.black
        )
        
        case_types = [
            "Employment Dispute Resolution", "Breach of Contract Claim",
            "Intellectual Property Violation", "Commercial Litigation Matter",
            "Partnership Dissolution Proceeding", "Regulatory Compliance Issue"
        ]
        
        case_type = random.choice(case_types)
        case_num = f"{random.randint(2023, 2024)}-{random.choice(['CV', 'EMP', 'IP', 'COM', 'REG'])}-{random.randint(1000, 9999)}"
        plaintiff = random.choice(self.companies)
        defendant = random.choice([c for c in self.companies if c != plaintiff])
        court_date = self.random_date(90)
        
        # Document header
        story.append(Paragraph("SUPERIOR COURT OF JUSTICE", legal_style))
        story.append(Paragraph("COMMERCIAL DIVISION", styles['Heading3']))
        story.append(Spacer(1, 30))
        
        # Case information
        case_info = [
            ['Case Number:', case_num],
            ['Case Type:', case_type],
            ['Filing Date:', court_date.strftime('%B %d, %Y')],
            ['Plaintiff:', plaintiff],
            ['Defendant:', defendant]
        ]
        
        case_table = Table(case_info, colWidths=[2*inch, 4*inch])
        case_table.setStyle(TableStyle([
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 11),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        story.append(case_table)
        story.append(Spacer(1, 30))
        
        # Legal notice content
        story.append(Paragraph("<b>NOTICE OF LEGAL PROCEEDINGS</b>", styles['Heading2']))
        story.append(Spacer(1, 15))
        
        legal_content = f"""
TO: {defendant}

TAKE NOTICE that a legal proceeding has been commenced against you by {plaintiff} 
for {case_type.lower()} and associated claims for damages and relief.

NATURE OF CLAIM:
The plaintiff alleges that the defendant has committed the following violations:
• Breach of contractual obligations and failure to perform agreed services
• Violation of confidentiality and non-disclosure agreements  
• Misappropriation of proprietary information and trade secrets
• Interference with business relationships and competitive practices

RELIEF SOUGHT:
The plaintiff seeks the following relief from this Court:
1. Monetary damages in the amount of {self.format_currency(random.randint(100000, 2000000))}
2. Injunctive relief to prevent further violations
3. Restitution of profits and gains wrongfully obtained
4. Pre-judgment and post-judgment interest
5. Attorney fees and costs of litigation
6. Such other relief as the Court deems just and proper

RESPONSE REQUIRED:
Any party wishing to defend this action must file a Statement of Defense within 
{random.choice([20, 30, 45])} days of service of this Notice. Failure to defend may 
result in judgment being granted against you without further notice.

For more information regarding this proceeding, contact the Court Registry or 
legal counsel at the address below.

DATED this {court_date.strftime('%d')} day of {court_date.strftime('%B')}, {court_date.year}.

                                    _________________________________
                                    Registrar, Superior Court of Justice
                                    
Attorney for Plaintiff:
{random.choice(['Thompson & Associates Legal LLP', 'Mitchell, Brown & Partners', 'Sterling Legal Group'])}
{random.choice(self.addresses)[0]}
{random.choice(self.addresses)[1]}, {random.choice(self.addresses)[2]} {random.choice(self.addresses)[3]}
Tel: {random.randint(555, 999)}-{random.randint(100, 999)}-{random.randint(1000, 9999)}
"""
        
        # Split content into paragraphs
        for paragraph in legal_content.strip().split('\n\n'):
            if paragraph.strip():
                story.append(Paragraph(paragraph.strip(), styles['Normal']))
                story.append(Spacer(1, 12))
        
        doc.build(story)
        return filepath

    def generate_report_docx(self, filename: str) -> str:
        """Generate a business report DOCX."""
        filepath = os.path.join(self.output_dir, 'docx', filename)
        doc = DocxDocument()
        
        # Title page
        title = doc.add_heading('QUARTERLY BUSINESS PERFORMANCE REPORT', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        quarter = random.choice(['Q1', 'Q2', 'Q3', 'Q4'])
        year = random.choice([2023, 2024])
        company = random.choice(self.companies)
        
        doc.add_paragraph(f"{quarter} {year} Executive Summary", style='Heading 1').alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f"{company}", style='Heading 2').alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("")
        
        # Executive Summary
        doc.add_heading('EXECUTIVE SUMMARY', level=1)
        
        revenue = random.randint(5000000, 50000000)
        growth = random.randint(-10, 35)
        profit_margin = random.uniform(8, 25)
        
        summary_text = f"""
This report presents the comprehensive business performance analysis for {quarter} {year}, 
highlighting key financial metrics, operational achievements, and strategic initiatives.

Our organization achieved total revenue of {self.format_currency(revenue)} during this quarter, 
representing a {growth}% {'increase' if growth > 0 else 'decrease'} compared to the same period 
last year. The profit margin improved to {profit_margin:.1f}%, demonstrating effective cost 
management and operational efficiency improvements.
"""
        doc.add_paragraph(summary_text.strip())
        
        # Financial Performance
        doc.add_heading('FINANCIAL PERFORMANCE', level=1)
        
        # Create a table for financial data
        table = doc.add_table(rows=1, cols=3)
        table.style = 'Table Grid'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Metric'
        hdr_cells[1].text = f'{quarter} {year}'
        hdr_cells[2].text = 'YoY Change'
        
        financial_data = [
            ('Total Revenue', self.format_currency(revenue), f"{growth:+.1f}%"),
            ('Operating Expenses', self.format_currency(revenue * 0.7), f"{random.randint(-5, 15):+.1f}%"),
            ('Net Profit', self.format_currency(revenue * profit_margin / 100), f"{random.randint(5, 25):+.1f}%"),
            ('Profit Margin', f"{profit_margin:.1f}%", f"{random.uniform(-2, 5):+.1f}pp")
        ]
        
        for metric, value, change in financial_data:
            row_cells = table.add_row().cells
            row_cells[0].text = metric
            row_cells[1].text = value
            row_cells[2].text = change
        
        # Operational Highlights
        doc.add_heading('OPERATIONAL HIGHLIGHTS', level=1)
        
        highlights = [
            f"Successfully launched {random.randint(2, 5)} new product initiatives",
            f"Expanded market presence in {random.randint(3, 8)} additional regions",
            f"Achieved customer satisfaction score of {random.uniform(4.2, 4.9):.1f}/5.0",
            f"Reduced operational costs by {random.uniform(3, 12):.1f}% through efficiency improvements",
            f"Increased employee retention rate to {random.uniform(85, 95):.1f}%"
        ]
        
        for highlight in highlights:
            doc.add_paragraph(highlight, style='List Bullet')
        
        # Recommendations
        doc.add_heading('STRATEGIC RECOMMENDATIONS', level=1)
        
        recommendations = [
            "Continue investment in high-performing product lines and market segments",
            "Implement advanced analytics tools to enhance decision-making capabilities", 
            "Strengthen customer relationship management and retention programs",
            "Explore strategic partnerships to accelerate market expansion",
            "Invest in employee development and training programs"
        ]
        
        for rec in recommendations:
            doc.add_paragraph(rec, style='List Number')
        
        # Conclusion
        doc.add_heading('CONCLUSION', level=1)
        conclusion = f"""
{quarter} {year} demonstrated {'strong' if growth > 10 else 'stable' if growth > 0 else 'challenging'} 
performance across key business metrics. The organization is well-positioned for continued growth 
and success in the upcoming quarters, with solid financial fundamentals and operational excellence 
driving sustainable competitive advantage.
"""
        doc.add_paragraph(conclusion.strip())
        
        doc.save(filepath)
        return filepath

    def generate_other_txt(self, filename: str) -> str:
        """Generate miscellaneous document TXT."""
        filepath = os.path.join(self.output_dir, 'txt', filename)
        
        doc_types = [
            "Technical Specification Document", "User Manual", "Meeting Minutes",
            "Project Status Update", "Training Material", "General Correspondence"
        ]
        
        doc_type = random.choice(doc_types)
        date = self.random_date(60)
        author_name, author_last, _ = random.choice(self.people)
        
        content = f"""
{doc_type.upper()}
{'='*len(doc_type.upper())}

Document Information:
- Title: {doc_type}
- Date: {date.strftime('%B %d, %Y')}
- Author: {author_name} {author_last}
- Version: {random.randint(1, 5)}.{random.randint(0, 9)}
- Reference: DOC-{random.randint(1000, 9999)}

{'='*60}
CONTENT OVERVIEW
{'='*60}

This document serves as {random.choice([
    'comprehensive reference material for various operational procedures',
    'detailed guidance for system implementation and configuration',
    'general information resource for project stakeholders',
    'supplementary documentation supporting business processes'
])}.

Key Topics Covered:
• {random.choice(['System requirements and specifications', 'Operational procedures and guidelines'])}
• {random.choice(['Implementation strategies and best practices', 'Process workflows and methodologies'])}
• {random.choice(['Quality assurance and testing protocols', 'Maintenance and support procedures'])}
• {random.choice(['Compliance requirements and standards', 'Performance metrics and evaluation criteria'])}

Technical Details:
{random.choice([
    'This section contains detailed technical specifications and configuration parameters required for proper system implementation.',
    'The following information provides comprehensive guidance for users and administrators regarding system operation.',
    'This documentation includes step-by-step procedures and troubleshooting guidelines for common scenarios.',
    'Reference material contained herein supports ongoing operations and maintenance activities.'
])}

Procedures and Guidelines:
1. Initial setup and configuration requirements
2. Standard operating procedures and workflows  
3. Quality control and validation processes
4. Documentation and reporting requirements
5. Escalation procedures and support contacts

Additional Information:
{random.choice([
    'Please refer to the appendices for detailed technical specifications and supplementary resources.',
    'For additional support or clarification, contact the appropriate department or project lead.',
    'This document should be reviewed and updated regularly to ensure accuracy and relevance.',
    'All users must familiarize themselves with the contents before proceeding with implementation.'
])}

Notes and Disclaimers:
- This document is provided for informational purposes only
- Procedures may vary based on specific system configurations
- Always follow current organizational policies and guidelines
- Contact technical support for assistance with complex scenarios

Document prepared by: {author_name} {author_last}
Department: {random.choice(['IT', 'Operations', 'Quality Assurance', 'Project Management'])}
Last Updated: {date.strftime('%B %d, %Y')}

{'='*60}
END OF DOCUMENT
{'='*60}
"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return filepath

    def generate_invoice_txt(self, filename: str) -> str:
        """Generate an invoice in TXT format."""
        filepath = os.path.join(self.output_dir, 'txt', filename)
        company = random.choice(self.companies)
        customer_name, customer_last, _ = random.choice(self.people)
        address = random.choice(self.addresses)
        invoice_date = self.random_date(90)
        due_date = invoice_date + timedelta(days=30)
        
        # Generate line items
        line_items = []
        subtotal = 0
        for i in range(random.randint(2, 5)):
            service = random.choice([
                'Software Development', 'Consulting Services', 'Project Management',
                'Technical Support', 'System Analysis', 'Training Services'
            ])
            hours = random.randint(10, 80)
            rate = random.choice([75, 85, 95, 105, 115, 125])
            amount = hours * rate
            subtotal += amount
            line_items.append(f"{service:<25} {hours:>6} hrs  ${rate:>6}/hr  ${amount:>8,.2f}")
        
        tax_rate = 0.0875
        tax_amount = subtotal * tax_rate
        total = subtotal + tax_amount
        
        content = f"""
================================================================================
                              INVOICE
================================================================================

{company}
{address}

Invoice To:                          Invoice Details:
{customer_name} {customer_last}      Invoice #: INV-{random.randint(10000, 99999)}
Business Client                      Date: {invoice_date.strftime('%B %d, %Y')}
                                    Due Date: {due_date.strftime('%B %d, %Y')}
                                    Terms: Net 30 Days

================================================================================
                            SERVICES PROVIDED
================================================================================

Description                    Hours    Rate      Amount
--------------------------------------------------------------------------------
{chr(10).join(line_items)}
                                              ________________
                                    Subtotal:  ${subtotal:>8,.2f}
                                        Tax:   ${tax_amount:>8,.2f}
                                              ================
                                       Total:  ${total:>8,.2f}

================================================================================
                              PAYMENT TERMS
================================================================================

Payment is due within 30 days of invoice date.
Please remit payment to the address above.
Thank you for your business!

Account Details:
- Account Number: {random.randint(100000, 999999)}
- Reference: INV-{random.randint(10000, 99999)}
- Customer ID: CUST-{random.randint(1000, 9999)}

Questions? Contact us at: billing@{company.lower().replace(' ', '').replace(',', '').replace('.', '')}.com

================================================================================
"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return filepath

    def generate_memo_txt(self, filename: str) -> str:
        """Generate a memo in TXT format."""
        filepath = os.path.join(self.output_dir, 'txt', filename)
        company = random.choice(self.companies)
        from_name, from_last, from_title = random.choice(self.people)
        to_name, to_last, to_title = random.choice(self.people)
        date = self.random_date(30)
        
        subjects = [
            'Policy Update - Remote Work Guidelines',
            'Quarterly Team Meeting Schedule',
            'New Safety Procedures Implementation',
            'Budget Allocation for Q4 Projects',
            'Performance Review Process Updates',
            'Holiday Schedule Announcement',
            'Training Program Requirements',
            'Office Relocation Timeline'
        ]
        
        subject = random.choice(subjects)
        
        content = f"""
================================================================================
                           INTERNAL MEMORANDUM
================================================================================

{company}

TO:      {to_name} {to_last}, {to_title}
FROM:    {from_name} {from_last}, {from_title}
DATE:    {date.strftime('%B %d, %Y')}
SUBJECT: {subject}

================================================================================

Dear Team,

{random.choice([
    'I am writing to inform you of important updates that will take effect immediately.',
    'Please be advised of the following policy changes and implementation timeline.',
    'This memo outlines new procedures that all staff members must follow.',
    'The following information requires your immediate attention and action.'
])}

Key Points:
• {random.choice(['All employees must complete the new training by month end', 'Department heads should schedule team meetings within two weeks'])}
• {random.choice(['Policy updates will be posted on the company intranet', 'Updated procedures are available in the employee handbook'])}
• {random.choice(['Questions should be directed to HR for clarification', 'Implementation begins on the first of next month'])}

Action Items:
1. Review the attached documentation carefully
2. {random.choice(['Schedule required training sessions for your team', 'Update departmental procedures accordingly'])}
3. {random.choice(['Submit compliance confirmation by the deadline', 'Coordinate with other departments as needed'])}
4. Follow up with any questions or concerns

Implementation Timeline:
- Phase 1: {random.choice(['Training completion', 'Policy review'])} - {(date + timedelta(days=14)).strftime('%B %d')}
- Phase 2: {random.choice(['Full implementation', 'Compliance verification'])} - {(date + timedelta(days=30)).strftime('%B %d')}

{random.choice([
    'Your cooperation and prompt attention to this matter is greatly appreciated.',
    'Please ensure all team members are informed of these important updates.',
    'Contact me directly if you need clarification on any of these points.',
    'Thank you for your continued dedication to our organizational excellence.'
])}

Best regards,

{from_name} {from_last}
{from_title}
{company}

================================================================================
"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return filepath

    def generate_legal_txt(self, filename: str) -> str:
        """Generate a legal document in TXT format."""
        filepath = os.path.join(self.output_dir, 'txt', filename)
        plaintiff_name, plaintiff_last, _ = random.choice(self.people)
        defendant_name, defendant_last, _ = random.choice(self.people)
        case_number = f"{random.randint(2020, 2024)}-CV-{random.randint(1000, 9999)}"
        date = self.random_date(60)
        
        content = f"""
================================================================================
                         SUPERIOR COURT OF JUSTICE
================================================================================

Case No: {case_number}

{plaintiff_name} {plaintiff_last}
                                                    Plaintiff
v.

{defendant_name} {defendant_last}
                                                    Defendant

================================================================================
                              LEGAL NOTICE
================================================================================

TO: {defendant_name} {defendant_last}

YOU ARE HEREBY NOTIFIED that a legal action has been commenced against you in the 
Superior Court of Justice. The nature of the claim against you is set forth in 
the Statement of Claim served with this Notice.

NATURE OF CLAIM:
{random.choice([
    'Breach of employment contract and violation of confidentiality agreement',
    'Business dispute regarding contractual obligations and damages',
    'Professional negligence and breach of fiduciary duty',
    'Violation of non-compete agreement and trade secret misappropriation'
])}

RELIEF SOUGHT:
The Plaintiff seeks monetary damages in the amount of ${random.randint(50000, 500000):,}, 
plus costs, interest, and such other relief as this Court deems just and proper.

YOU HAVE THIRTY (30) DAYS after service of this Notice to file an Answer with 
the Court. Failure to respond within the prescribed time may result in a 
default judgment being entered against you.

COURT INFORMATION:
Superior Court of Justice
Civil Division
{random.choice(self.addresses)}

ATTORNEY FOR PLAINTIFF:
{random.choice(['Law Offices of Smith & Associates', 'Johnson Legal Group', 'Williams & Partners LLP'])}
Attorney Bar No: {random.randint(100000, 999999)}
{random.choice(self.addresses)}

DATE OF SERVICE: {date.strftime('%B %d, %Y')}

IMPORTANT: If you fail to file an Answer within thirty (30) days after service 
of this Notice, judgment by default may be taken against you for the relief 
demanded in the Statement of Claim.

YOU ARE ADVISED TO SEEK LEGAL COUNSEL immediately if you have not already done so.

================================================================================

CERTIFICATE OF SERVICE

I hereby certify that a true and correct copy of the foregoing Notice was served 
upon the defendant by the following method:

[ ] Personal service
[ ] Certified mail, return receipt requested
[ ] Publication in newspaper of general circulation

Date: {date.strftime('%B %d, %Y')}

                                    _________________________
                                    Clerk of Court

================================================================================
"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return filepath

    def generate_report_pdf(self, filename: str) -> str:
        """Generate a report in PDF format."""
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        
        filepath = os.path.join(self.output_dir, 'pdf', filename)
        doc = SimpleDocTemplate(filepath, pagesize=letter, topMargin=0.5*inch)
        
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=18,
            textColor=colors.darkblue,
            spaceAfter=20
        )
        
        company = random.choice(self.companies)
        quarter = random.choice(['Q1', 'Q2', 'Q3', 'Q4'])
        year = random.choice([2023, 2024])
        date = self.random_date(60)
        
        elements = []
        
        # Title
        title = Paragraph(f"{quarter} {year} Business Performance Report", title_style)
        elements.append(title)
        elements.append(Spacer(1, 0.2*inch))
        
        # Company info
        company_info = Paragraph(f"<b>{company}</b><br/>Quarterly Business Analysis", styles['Normal'])
        elements.append(company_info)
        elements.append(Spacer(1, 0.3*inch))
        
        # Executive Summary
        exec_summary = Paragraph("<b>EXECUTIVE SUMMARY</b>", styles['Heading2'])
        elements.append(exec_summary)
        
        summary_text = f"""
        This report provides a comprehensive analysis of business performance for {quarter} {year}. 
        Key metrics show {random.choice(['strong growth', 'steady improvement', 'positive trends'])} across 
        multiple operational areas. Revenue increased by {random.randint(5, 25)}% compared to the previous quarter, 
        with {random.choice(['customer satisfaction', 'operational efficiency', 'market penetration'])} reaching 
        new highs.
        """
        elements.append(Paragraph(summary_text, styles['Normal']))
        elements.append(Spacer(1, 0.2*inch))
        
        # Financial Performance Table
        financial_header = Paragraph("<b>FINANCIAL PERFORMANCE</b>", styles['Heading2'])
        elements.append(financial_header)
        
        financial_data = [
            ['Metric', 'Current Quarter', 'Previous Quarter', 'Change'],
            ['Revenue', f'${random.randint(500, 2000)}K', f'${random.randint(400, 1800)}K', f'+{random.randint(5, 25)}%'],
            ['Profit Margin', f'{random.randint(15, 35)}%', f'{random.randint(10, 30)}%', f'+{random.randint(2, 8)}%'],
            ['Operating Costs', f'${random.randint(300, 1500)}K', f'${random.randint(350, 1600)}K', f'-{random.randint(5, 15)}%'],
            ['Customer Acquisition', f'{random.randint(150, 500)}', f'{random.randint(100, 400)}', f'+{random.randint(10, 40)}%']
        ]
        
        financial_table = Table(financial_data, colWidths=[2*inch, 1.5*inch, 1.5*inch, 1*inch])
        financial_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.lightblue),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.darkblue),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.lightgrey),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        elements.append(financial_table)
        elements.append(Spacer(1, 0.3*inch))
        
        # Key Achievements
        achievements = Paragraph("<b>KEY ACHIEVEMENTS</b>", styles['Heading2'])
        elements.append(achievements)
        
        achievement_text = f"""
        • Exceeded revenue targets by {random.randint(8, 20)}%<br/>
        • Improved customer retention rate to {random.randint(85, 95)}%<br/>
        • Launched {random.randint(2, 5)} new {random.choice(['products', 'services', 'initiatives'])}<br/>
        • Reduced operational costs by {random.randint(10, 25)}%<br/>
        • Achieved {random.randint(90, 98)}% customer satisfaction rating
        """
        elements.append(Paragraph(achievement_text, styles['Normal']))
        
        doc.build(elements)
        return filepath

    def generate_contract_pdf(self, filename: str) -> str:
        """Generate a contract in PDF format."""
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        
        filepath = os.path.join(self.output_dir, 'pdf', filename)
        doc = SimpleDocTemplate(filepath, pagesize=letter, topMargin=0.5*inch)
        
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=16,
            textColor=colors.darkblue,
            spaceAfter=20
        )
        
        company1 = random.choice(self.companies)
        company2 = random.choice([c for c in self.companies if c != company1])
        party1_name, party1_last, party1_title = random.choice(self.people)
        party2_name, party2_last, party2_title = random.choice(self.people)
        date = self.random_date(30)
        
        elements = []
        
        # Title
        title = Paragraph("PROFESSIONAL SERVICES AGREEMENT", title_style)
        elements.append(title)
        elements.append(Spacer(1, 0.3*inch))
        
        # Agreement text
        agreement_text = f"""
        This Professional Services Agreement ("Agreement") is entered into on {date.strftime('%B %d, %Y')} 
        between {company1} ("Provider") and {company2} ("Client").
        
        <b>WHEREAS</b>, Provider possesses expertise in professional consulting services; and
        
        <b>WHEREAS</b>, Client desires to engage Provider for such services;
        
        <b>NOW THEREFORE</b>, the parties agree as follows:
        
        <b>1. SERVICES</b><br/>
        Provider shall provide {random.choice(['software development', 'consulting', 'project management', 'technical support'])} 
        services as detailed in Exhibit A, which is incorporated herein by reference.
        
        <b>2. COMPENSATION</b><br/>
        Client shall pay Provider ${random.randint(50, 200):,} per month for the duration of this Agreement. 
        Payment is due within 30 days of invoice receipt.
        
        <b>3. TERM</b><br/>
        This Agreement shall commence on {date.strftime('%B %d, %Y')} and continue for a period of 
        {random.choice(['twelve (12)', 'twenty-four (24)', 'thirty-six (36)'])} months, unless terminated earlier.
        
        <b>4. CONFIDENTIALITY</b><br/>
        Both parties acknowledge that confidential information may be shared and agree to maintain 
        strict confidentiality of all proprietary information.
        
        <b>5. TERMINATION</b><br/>
        Either party may terminate this Agreement with {random.choice(['30', '60', '90'])} days written notice.
        
        <b>6. GOVERNING LAW</b><br/>
        This Agreement shall be governed by the laws of the State of {random.choice(['California', 'New York', 'Texas'])}.
        
        <b>IN WITNESS WHEREOF</b>, the parties have executed this Agreement as of the date first 
        written above.
        """
        
        elements.append(Paragraph(agreement_text, styles['Normal']))
        elements.append(Spacer(1, 0.5*inch))
        
        # Signature blocks
        signature_text = f"""
        <b>PROVIDER:</b><br/>
        {company1}<br/><br/>
        By: _________________________<br/>
        {party1_name} {party1_last}, {party1_title}<br/>
        Date: ____________________<br/><br/>
        
        <b>CLIENT:</b><br/>
        {company2}<br/><br/>
        By: _________________________<br/>
        {party2_name} {party2_last}, {party2_title}<br/>
        Date: ____________________
        """
        
        elements.append(Paragraph(signature_text, styles['Normal']))
        
        doc.build(elements)
        return filepath

    def generate_contract_docx(self, filename: str) -> str:
        """Generate a contract in DOCX format."""
        from docx import Document
        from docx.shared import Inches
        
        filepath = os.path.join(self.output_dir, 'docx', filename)
        doc = Document()
        
        company1 = random.choice(self.companies)
        company2 = random.choice([c for c in self.companies if c != company1])
        party1_name, party1_last, party1_title = random.choice(self.people)
        party2_name, party2_last, party2_title = random.choice(self.people)
        date = self.random_date(30)
        
        # Title
        title = doc.add_heading('EMPLOYMENT CONTRACT', 0)
        title.alignment = 1  # Center alignment
        
        # Header
        doc.add_paragraph()
        header = doc.add_paragraph()
        header.add_run('Contract Number: ').bold = True
        header.add_run(f'EMP-{random.randint(10000, 99999)}')
        header.add_run('\nEffective Date: ').bold = True
        header.add_run(date.strftime('%B %d, %Y'))
        
        doc.add_paragraph()
        
        # Parties
        parties_para = doc.add_paragraph()
        parties_para.add_run('AGREEMENT BETWEEN:\n').bold = True
        parties_para.add_run(f'Employer: {company1}\n')
        parties_para.add_run(f'Employee: {party1_name} {party1_last}\n')
        
        doc.add_paragraph()
        
        # Terms
        terms = doc.add_heading('TERMS AND CONDITIONS', level=1)
        
        # Section 1
        section1 = doc.add_paragraph()
        section1.add_run('1. POSITION AND DUTIES\n').bold = True
        section1.add_run(f'Employee shall serve as {party1_title} and shall perform duties as assigned by the Employer. Employee agrees to devote full professional time and attention to the business of the Employer.')
        
        doc.add_paragraph()
        
        # Section 2
        section2 = doc.add_paragraph()
        section2.add_run('2. COMPENSATION\n').bold = True
        salary = random.randint(60000, 150000)
        section2.add_run(f'Employee shall receive an annual salary of ${salary:,}, payable in accordance with Employer\'s standard payroll practices. Employee shall be eligible for annual performance bonuses at Employer\'s discretion.')
        
        doc.add_paragraph()
        
        # Section 3
        section3 = doc.add_paragraph()
        section3.add_run('3. BENEFITS\n').bold = True
        section3.add_run('Employee shall be entitled to participate in all employee benefit plans and programs made available to employees of similar rank and tenure, subject to the terms and conditions of such plans.')
        
        doc.add_paragraph()
        
        # Section 4
        section4 = doc.add_paragraph()
        section4.add_run('4. CONFIDENTIALITY\n').bold = True
        section4.add_run('Employee acknowledges that during employment, Employee will have access to confidential information. Employee agrees to maintain strict confidentiality of all proprietary information.')
        
        doc.add_paragraph()
        
        # Section 5
        section5 = doc.add_paragraph()
        section5.add_run('5. TERMINATION\n').bold = True
        section5.add_run(f'This agreement may be terminated by either party with {random.choice(["30", "60", "90"])} days written notice. Upon termination, Employee agrees to return all company property.')
        
        doc.add_paragraph()
        
        # Signatures
        doc.add_paragraph()
        sig_para = doc.add_paragraph()
        sig_para.add_run('SIGNATURES:\n').bold = True
        sig_para.add_run(f'\nEmployer: {company1}\n\n')
        sig_para.add_run(f'Employee: {party1_name} {party1_last}\n\n')
        sig_para.add_run('Signature: _________________________  Date: ___________\n\n')
        sig_para.add_run(f'Employee: {party1_name} {party1_last}\n\n')
        sig_para.add_run('Signature: _________________________  Date: ___________')
        
        doc.save(filepath)
        return filepath

    def generate_other_pdf(self, filename: str) -> str:
        """Generate an 'other' document in PDF format."""
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        
        filepath = os.path.join(self.output_dir, 'pdf', filename)
        doc = SimpleDocTemplate(filepath, pagesize=letter, topMargin=0.5*inch)
        
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=16,
            textColor=colors.darkblue,
            spaceAfter=20
        )
        
        doc_types = ['Technical Specification', 'User Manual', 'Project Proposal', 'Meeting Minutes']
        doc_type = random.choice(doc_types)
        company = random.choice(self.companies)
        author_name, author_last, author_title = random.choice(self.people)
        date = self.random_date(60)
        
        elements = []
        
        # Title
        title = Paragraph(f"{doc_type.upper()}", title_style)
        elements.append(title)
        elements.append(Spacer(1, 0.2*inch))
        
        # Header info
        header_text = f"""
        <b>Document Information:</b><br/>
        Company: {company}<br/>
        Prepared by: {author_name} {author_last}, {author_title}<br/>
        Date: {date.strftime('%B %d, %Y')}<br/>
        Version: {random.randint(1, 5)}.{random.randint(0, 9)}<br/>
        Reference: DOC-{random.randint(1000, 9999)}
        """
        elements.append(Paragraph(header_text, styles['Normal']))
        elements.append(Spacer(1, 0.3*inch))
        
        # Content
        if doc_type == 'Technical Specification':
            content_text = f"""
            <b>OVERVIEW</b><br/>
            This document outlines the technical specifications for system implementation. 
            The solution addresses {random.choice(['scalability requirements', 'performance optimization', 'security compliance', 'integration needs'])} 
            while maintaining compatibility with existing infrastructure.
            
            <b>REQUIREMENTS</b><br/>
            • Minimum system requirements: {random.choice(['8GB RAM', '16GB RAM', '32GB RAM'])}, 
            {random.choice(['2.5GHz processor', '3.0GHz processor', '3.5GHz processor'])}<br/>
            • Software dependencies: {random.choice(['Java 11+', 'Python 3.8+', '.NET 6.0+'])}<br/>
            • Database: {random.choice(['PostgreSQL 13+', 'MySQL 8.0+', 'SQL Server 2019+'])}<br/>
            • Network: {random.choice(['1Gbps', '10Gbps', '100Mbps'])} minimum bandwidth
            
            <b>IMPLEMENTATION NOTES</b><br/>
            Installation should be performed during maintenance windows. 
            All configurations must be validated before production deployment.
            """
        else:
            content_text = f"""
            <b>SUMMARY</b><br/>
            This document provides {random.choice(['comprehensive guidance', 'detailed procedures', 'reference information'])} 
            for {random.choice(['project stakeholders', 'system users', 'technical teams', 'business analysts'])}.
            
            <b>KEY TOPICS</b><br/>
            • {random.choice(['System configuration and setup procedures', 'User interface and navigation guidelines'])}<br/>
            • {random.choice(['Troubleshooting and maintenance protocols', 'Best practices and recommendations'])}<br/>
            • {random.choice(['Security considerations and compliance requirements', 'Performance optimization techniques'])}
            
            <b>ADDITIONAL INFORMATION</b><br/>
            For technical support or additional questions, please contact the 
            {random.choice(['IT Help Desk', 'Project Team', 'System Administrator', 'Technical Lead'])} 
            at extension {random.randint(1000, 9999)}.
            """
        
        elements.append(Paragraph(content_text, styles['Normal']))
        
        doc.build(elements)
        return filepath

    def generate_other_docx(self, filename: str) -> str:
        """Generate an 'other' document in DOCX format."""
        from docx import Document
        
        filepath = os.path.join(self.output_dir, 'docx', filename)
        doc = Document()
        
        doc_types = ['Meeting Minutes', 'Project Proposal', 'Technical Manual', 'Policy Document']
        doc_type = random.choice(doc_types)
        company = random.choice(self.companies)
        author_name, author_last, author_title = random.choice(self.people)
        date = self.random_date(60)
        
        # Title
        title = doc.add_heading(f'{doc_type.upper()}', 0)
        title.alignment = 1
        
        # Header
        doc.add_paragraph()
        header = doc.add_paragraph()
        header.add_run('Company: ').bold = True
        header.add_run(f'{company}\n')
        header.add_run('Prepared by: ').bold = True
        header.add_run(f'{author_name} {author_last}, {author_title}\n')
        header.add_run('Date: ').bold = True
        header.add_run(date.strftime('%B %d, %Y'))
        
        doc.add_paragraph()
        
        # Content based on document type
        if doc_type == 'Meeting Minutes':
            # Meeting details
            meeting_para = doc.add_paragraph()
            meeting_para.add_run('MEETING DETAILS\n').bold = True
            meeting_para.add_run(f'Date: {date.strftime("%B %d, %Y")}\n')
            meeting_para.add_run(f'Time: {random.choice(["9:00 AM", "10:00 AM", "2:00 PM", "3:00 PM"])}\n')
            meeting_para.add_run(f'Location: {random.choice(["Conference Room A", "Virtual Meeting", "Executive Boardroom"])}\n')
            
            # Attendees
            attendees_para = doc.add_paragraph()
            attendees_para.add_run('ATTENDEES\n').bold = True
            for i in range(random.randint(3, 6)):
                name, last, title = random.choice(self.people)
                attendees_para.add_run(f'• {name} {last}, {title}\n')
            
            # Discussion points
            discussion_para = doc.add_paragraph()
            discussion_para.add_run('DISCUSSION POINTS\n').bold = True
            discussion_para.add_run('• Quarterly budget review and allocation planning\n')
            discussion_para.add_run('• Project timeline updates and milestone tracking\n')
            discussion_para.add_run('• Resource allocation and staffing requirements\n')
            discussion_para.add_run('• Process improvement initiatives and implementation\n')
            
            # Action items
            action_para = doc.add_paragraph()
            action_para.add_run('ACTION ITEMS\n').bold = True
            action_para.add_run(f'• Complete budget analysis by {(date + timedelta(days=7)).strftime("%B %d")}\n')
            action_para.add_run(f'• Schedule follow-up meetings with department heads\n')
            action_para.add_run(f'• Review and update project documentation\n')
        
        else:
            # General content
            overview_para = doc.add_paragraph()
            overview_para.add_run('OVERVIEW\n').bold = True
            overview_para.add_run(f'This {doc_type.lower()} provides {random.choice(["comprehensive information", "detailed guidance", "reference material"])} for {random.choice(["project stakeholders", "team members", "system users"])}. ')
            overview_para.add_run(f'The content addresses {random.choice(["operational requirements", "implementation procedures", "best practices", "compliance standards"])} and related considerations.')
            
            doc.add_paragraph()
            
            # Key sections
            sections_para = doc.add_paragraph()
            sections_para.add_run('KEY SECTIONS\n').bold = True
            sections_para.add_run('• System requirements and configuration guidelines\n')
            sections_para.add_run('• Implementation procedures and best practices\n')
            sections_para.add_run('• Quality assurance and testing protocols\n')
            sections_para.add_run('• Maintenance and support procedures\n')
            
            doc.add_paragraph()
            
            # Notes
            notes_para = doc.add_paragraph()
            notes_para.add_run('IMPORTANT NOTES\n').bold = True
            notes_para.add_run('Please review all sections carefully before proceeding with implementation. ')
            notes_para.add_run('Contact the project team for questions or clarification on specific procedures. ')
            notes_para.add_run('This document should be updated regularly to reflect current practices and requirements.')
        
        doc.save(filepath)
        return filepath

    def create_memo_content(self, fmt: str) -> dict:
        """Create enhanced memo content."""
        company = random.choice(self.companies)
        sender_name, sender_last, sender_title = random.choice(self.people)
        subject = random.choice(self.memo_subjects)
        date = self.random_date(30).strftime('%B %d, %Y')
        
        if fmt in ['html', 'rtf', 'odt']:
            body = f"""
            <div class="section">
                <h3>TO:</h3> <p>All Staff Members</p>
                <h3>FROM:</h3> <p>{sender_name} {sender_last}, {sender_title}</p>
                <h3>SUBJECT:</h3> <p>{subject}</p>
            </div>
            
            <div class="section">
                <h3>MEMO CONTENT:</h3>
                <p>This memorandum serves to inform all staff members of important updates and changes to company policy.</p>
                <p><strong>Effective Date:</strong> {(datetime.now() + timedelta(days=random.randint(7, 30))).strftime('%B %d, %Y')}</p>
                <p><strong>Implementation Details:</strong></p>
                <ul>
                    <li>All employees must review the updated procedures by the effective date</li>
                    <li>Department heads should schedule team meetings to discuss changes</li>
                    <li>Questions should be directed to Human Resources or your direct supervisor</li>
                    <li>Compliance with new policies is mandatory for all staff members</li>
                </ul>
                <p>Please ensure you understand these changes and contact HR if you have any questions.</p>
            </div>
            """
        else:
            body = f"""TO: All Staff Members
FROM: {sender_name} {sender_last}, {sender_title}
SUBJECT: {subject}

This memorandum serves to inform all staff members of important updates and changes to company policy.

Effective Date: {(datetime.now() + timedelta(days=random.randint(7, 30))).strftime('%B %d, %Y')}

IMPLEMENTATION DETAILS:
• All employees must review the updated procedures by the effective date
• Department heads should schedule team meetings to discuss changes  
• Questions should be directed to Human Resources or your direct supervisor
• Compliance with new policies is mandatory for all staff members

Please ensure you understand these changes and contact HR if you have any questions.

Best regards,
{sender_name} {sender_last}
{sender_title}"""

        return {
            'title': 'INTERNAL MEMORANDUM',
            'company': company,
            'date': date,
            'body': body,
            'author': f"{sender_name} {sender_last}"
        }

    def create_contract_content(self, fmt: str) -> dict:
        """Create enhanced contract content."""
        company1 = random.choice(self.companies)
        company2 = random.choice([c for c in self.companies if c != company1])
        contract_type = random.choice(self.contract_types)
        date = self.random_date(60).strftime('%B %d, %Y')
        
        if fmt in ['html', 'rtf', 'odt']:
            body = f"""
            <div class="section">
                <h3>PARTIES</h3>
                <p><strong>Provider:</strong> {company1}</p>
                <p><strong>Client:</strong> {company2}</p>
            </div>
            
            <div class="section">
                <h3>TERMS AND CONDITIONS</h3>
                <p>This agreement establishes the terms under which professional services will be provided.</p>
                <p><strong>Effective Date:</strong> {date}</p>
                <p><strong>Duration:</strong> {random.randint(12, 36)} months</p>
                <p><strong>Compensation:</strong> ${random.randint(50000, 200000):,} annually</p>
            </div>
            
            <div class="section">
                <h3>OBLIGATIONS</h3>
                <ul>
                    <li>Provider shall deliver services in accordance with industry standards</li>
                    <li>Client shall provide necessary access and cooperation</li>
                    <li>Both parties agree to maintain confidentiality of proprietary information</li>
                    <li>All deliverables become property of the client upon payment</li>
                </ul>
            </div>
            
            <div class="section">
                <h3>SIGNATURES</h3>
                <p>By signing below, both parties agree to the terms of this contract.</p>
                <p>Provider: _________________________ Date: _____________</p>
                <p>Client: _________________________ Date: _____________</p>
            </div>
            """
        else:
            body = f"""{contract_type.upper()}

PARTIES:
Provider: {company1}
Client: {company2}

TERMS AND CONDITIONS:
This agreement establishes the terms under which professional services will be provided.

Effective Date: {date}
Duration: {random.randint(12, 36)} months  
Compensation: ${random.randint(50000, 200000):,} annually

OBLIGATIONS:
• Provider shall deliver services in accordance with industry standards
• Client shall provide necessary access and cooperation
• Both parties agree to maintain confidentiality of proprietary information
• All deliverables become property of the client upon payment

GOVERNING LAW:
This Agreement shall be governed by the laws of the State of {random.choice(['California', 'New York', 'Texas'])}.

SIGNATURES:
By signing below, both parties agree to the terms of this contract.

Provider: _________________________ Date: _____________

Client: _________________________ Date: _____________"""

        return {
            'title': contract_type.upper(),
            'company': company1,
            'date': date,
            'body': body,
            'author': "Legal Department"
        }

    def create_legal_content(self, fmt: str) -> dict:
        """Create enhanced legal content."""
        court = f"Superior Court of {random.choice(['Justice', 'California', 'New York', 'Texas'])}"
        case_num = f"Case No. {random.randint(2023, 2024)}-{random.randint(100, 999)}"
        plaintiff_name, plaintiff_last, _ = random.choice(self.people)
        defendant_name, defendant_last, _ = random.choice(self.people)
        date = self.random_date(60).strftime('%B %d, %Y')
        
        if fmt in ['html', 'rtf', 'odt']:
            body = f"""
            <div class="section">
                <h3>CASE INFORMATION</h3>
                <p><strong>{case_num}</strong></p>
                <p><strong>Plaintiff:</strong> {plaintiff_name} {plaintiff_last}</p>
                <p><strong>Defendant:</strong> {defendant_name} {defendant_last}</p>
            </div>
            
            <div class="section">
                <h3>LEGAL NOTICE</h3>
                <p>YOU ARE HEREBY NOTIFIED that a legal action has been commenced against you in the {court}.</p>
                <p><strong>Nature of Claim:</strong> {random.choice(['Breach of contract and damages', 'Employment dispute', 'Business litigation'])}</p>
                <p><strong>Relief Sought:</strong> Monetary damages of ${random.randint(50000, 500000):,} plus costs and attorney fees.</p>
            </div>
            
            <div class="section">
                <h3>RESPONSE REQUIRED</h3>
                <p>You have <strong>THIRTY (30) DAYS</strong> after service of this notice to file an answer with the court.</p>
                <p>Failure to respond may result in a default judgment being entered against you.</p>
            </div>
            
            <div class="section">
                <h3>COURT INFORMATION</h3>
                <p>{court}<br/>
                Civil Division<br/>
                {random.choice(self.addresses)[0]}</p>
            </div>
            """
        else:
            body = f"""{court}
{case_num}

{plaintiff_name} {plaintiff_last}
                                                    Plaintiff
v.

{defendant_name} {defendant_last}
                                                    Defendant

LEGAL NOTICE

TO: {defendant_name} {defendant_last}

YOU ARE HEREBY NOTIFIED that a legal action has been commenced against you in the {court}. 

NATURE OF CLAIM:
{random.choice(['Breach of contract and violation of agreement terms', 'Employment dispute regarding contractual obligations', 'Business litigation and damages'])}

RELIEF SOUGHT:
The Plaintiff seeks monetary damages in the amount of ${random.randint(50000, 500000):,}, plus costs, interest, and attorney fees.

RESPONSE REQUIRED:
YOU HAVE THIRTY (30) DAYS after service of this Notice to file an Answer with the Court. Failure to respond within the prescribed time may result in a default judgment being entered against you.

COURT INFORMATION:
{court}
Civil Division
{random.choice(self.addresses)[0]}

ATTORNEY FOR PLAINTIFF:
{random.choice(['Law Offices of Smith & Associates', 'Johnson Legal Group', 'Williams & Partners LLP'])}
Attorney Bar No: {random.randint(100000, 999999)}"""

        return {
            'title': 'LEGAL NOTICE',
            'company': court,
            'date': date,
            'body': body,
            'author': "Court Clerk"
        }

    def create_report_content(self, fmt: str) -> dict:
        """Create enhanced report content."""
        company = random.choice(self.companies)
        author_name, author_last, author_title = random.choice(self.people)
        quarter = random.choice(['Q1', 'Q2', 'Q3', 'Q4'])
        year = random.choice([2023, 2024])
        date = self.random_date(30).strftime('%B %d, %Y')
        
        revenue = random.randint(500000, 2000000)
        expenses = random.randint(300000, int(revenue * 0.8))
        profit = revenue - expenses
        growth = random.uniform(-10, 25)
        
        if fmt in ['html', 'rtf', 'odt']:
            body = f"""
            <div class="section">
                <h3>EXECUTIVE SUMMARY</h3>
                <p>This report provides a comprehensive analysis of business performance for {quarter} {year}.</p>
                <p>Key highlights include revenue of ${revenue:,}, representing {growth:+.1f}% growth year-over-year.</p>
            </div>
            
            <div class="section">
                <h3>FINANCIAL PERFORMANCE</h3>
                <table>
                    <tr><th>Metric</th><th>Current Quarter</th><th>Previous Quarter</th><th>Change</th></tr>
                    <tr><td>Revenue</td><td>${revenue:,}</td><td>${int(revenue * 0.9):,}</td><td>+{int((revenue - revenue * 0.9) / (revenue * 0.9) * 100)}%</td></tr>
                    <tr><td>Expenses</td><td>${expenses:,}</td><td>${int(expenses * 1.1):,}</td><td>-{int((expenses * 1.1 - expenses) / (expenses * 1.1) * 100)}%</td></tr>
                    <tr><td>Net Profit</td><td>${profit:,}</td><td>${int(profit * 0.8):,}</td><td>+{int((profit - profit * 0.8) / (profit * 0.8) * 100)}%</td></tr>
                </table>
            </div>
            
            <div class="section">
                <h3>KEY METRICS</h3>
                <ul>
                    <li>Customer Satisfaction: {random.randint(85, 98)}%</li>
                    <li>Employee Retention: {random.randint(88, 96)}%</li>
                    <li>Market Share: {random.randint(12, 25)}%</li>
                    <li>Operational Efficiency: {random.randint(78, 92)}%</li>
                </ul>
            </div>
            
            <div class="section">
                <h3>RECOMMENDATIONS</h3>
                <p>Based on the analysis, we recommend continued investment in core business areas and strategic expansion into emerging markets.</p>
            </div>
            """
        else:
            body = f"""EXECUTIVE SUMMARY:
This report provides a comprehensive analysis of business performance for {quarter} {year}.
Key highlights include revenue of ${revenue:,}, representing {growth:+.1f}% growth year-over-year.

FINANCIAL PERFORMANCE:
• Revenue: ${revenue:,}
• Expenses: ${expenses:,}  
• Net Profit: ${profit:,}
• Profit Margin: {(profit/revenue*100):.1f}%

KEY METRICS:
• Customer Satisfaction: {random.randint(85, 98)}%
• Employee Retention: {random.randint(88, 96)}%
• Market Share: {random.randint(12, 25)}%
• Operational Efficiency: {random.randint(78, 92)}%

MARKET ANALYSIS:
The business environment showed {random.choice(['strong growth', 'steady performance', 'moderate challenges'])} 
during this quarter. Customer acquisition increased by {random.randint(5, 25)}% while retention 
rates remained stable at {random.randint(88, 95)}%.

RECOMMENDATIONS:
• Continue investment in core business areas
• Expand into emerging market segments  
• Optimize operational processes for efficiency
• Strengthen customer relationship management
• Invest in employee development and retention

Prepared by: {author_name} {author_last}, {author_title}"""

        return {
            'title': f'{quarter} {year} BUSINESS PERFORMANCE REPORT',
            'company': company,
            'date': date,
            'body': body,
            'author': f"{author_name} {author_last}"
        }

    def create_other_content(self, fmt: str) -> dict:
        """Create enhanced other document content."""
        company = random.choice(self.companies)
        author_name, author_last, author_title = random.choice(self.people)
        doc_types = ['Technical Manual', 'User Guide', 'System Documentation', 'Process Manual', 'Reference Guide']
        doc_type = random.choice(doc_types)
        date = self.random_date(60).strftime('%B %d, %Y')
        
        if fmt in ['html', 'rtf', 'odt']:
            body = f"""
            <div class="section">
                <h3>DOCUMENT OVERVIEW</h3>
                <p>This {doc_type.lower()} provides comprehensive guidance for system users and administrators.</p>
                <p><strong>Version:</strong> {random.randint(1, 5)}.{random.randint(0, 9)}</p>
                <p><strong>Last Updated:</strong> {date}</p>
            </div>
            
            <div class="section">
                <h3>SYSTEM REQUIREMENTS</h3>
                <ul>
                    <li>Operating System: {random.choice(['Windows 10+', 'macOS 12+', 'Linux Ubuntu 20.04+'])}</li>
                    <li>Memory: {random.choice(['8GB', '16GB', '32GB'])} RAM minimum</li>
                    <li>Storage: {random.choice(['50GB', '100GB', '250GB'])} available space</li>
                    <li>Network: Broadband internet connection required</li>
                </ul>
            </div>
            
            <div class="section">
                <h3>INSTALLATION PROCEDURES</h3>
                <ol>
                    <li>Download the installation package from the official website</li>
                    <li>Run the installer with administrator privileges</li>
                    <li>Follow the setup wizard instructions</li>
                    <li>Configure initial settings and user preferences</li>
                    <li>Complete the installation and restart the system</li>
                </ol>
            </div>
            
            <div class="section">
                <h3>TROUBLESHOOTING</h3>
                <p>For technical support, contact the IT Help Desk at extension {random.randint(1000, 9999)} or email support@{company.lower().replace(' ', '').replace(',', '').replace('.', '')}.com</p>
            </div>
            """
        else:
            body = f"""DOCUMENT OVERVIEW:
This {doc_type.lower()} provides comprehensive guidance for system users and administrators.

Version: {random.randint(1, 5)}.{random.randint(0, 9)}
Last Updated: {date}

SYSTEM REQUIREMENTS:
• Operating System: {random.choice(['Windows 10+', 'macOS 12+', 'Linux Ubuntu 20.04+'])}
• Memory: {random.choice(['8GB', '16GB', '32GB'])} RAM minimum
• Storage: {random.choice(['50GB', '100GB', '250GB'])} available space
• Network: Broadband internet connection required

INSTALLATION PROCEDURES:
1. Download the installation package from the official website
2. Run the installer with administrator privileges  
3. Follow the setup wizard instructions
4. Configure initial settings and user preferences
5. Complete the installation and restart the system

CONFIGURATION GUIDELINES:
• Set up user accounts and permissions
• Configure network settings and security protocols
• Customize user interface preferences
• Enable automatic updates and backup procedures

MAINTENANCE PROCEDURES:
• Perform regular system updates and patches
• Monitor system performance and resource usage
• Backup important data and configuration files
• Review security logs and access controls

TROUBLESHOOTING:
For technical support, contact the IT Help Desk at extension {random.randint(1000, 9999)} 
or email support@{company.lower().replace(' ', '').replace(',', '').replace('.', '')}.com

Quality Assurance: All procedures have been tested and validated.
Documentation Standards: This document follows ISO 9001 guidelines."""

        return {
            'title': doc_type.upper(),
            'company': company,
            'date': date,
            'body': body,
            'author': f"{author_name} {author_last}"
        }

    # ===========================================
    # ENHANCED MAIN GENERATION FUNCTION
    # ===========================================

    def generate_all_diverse_documents(self, total_docs: int = 120) -> Dict[str, int]:
        """Generate diverse documents across ALL supported formats and categories."""
        print("🚀 Generating diverse sample documents in MULTIPLE formats...")
        print(f"📊 Creating {total_docs} documents across 6 categories and 11 file formats")
        
        self.create_output_dir()
        
        # Enhanced distribution
        docs_per_category = total_docs // 6
        remainder = total_docs % 6
        
        distribution = {
            'invoice': docs_per_category + (1 if remainder > 0 else 0),
            'memo': docs_per_category + (1 if remainder > 1 else 0),
            'contract': docs_per_category + (1 if remainder > 2 else 0),
            'legal': docs_per_category + (1 if remainder > 3 else 0),
            'report': docs_per_category + (1 if remainder > 4 else 0),
            'other': docs_per_category + (1 if remainder > 5 else 0)
        }
        
        generated_files = []
        format_count = {'pdf': 0, 'docx': 0, 'txt': 0, 'doc': 0, 'html': 0, 'rtf': 0, 
                       'odt': 0, 'pptx': 0, 'ppt': 0, 'xlsx': 0, 'xls': 0}
        
        # Enhanced format distributions for each category
        category_formats = {
            'invoice': {
                'formats': ['pdf', 'xlsx', 'xls', 'html', 'docx', 'txt'],
                'weights': [0.4, 0.25, 0.15, 0.1, 0.07, 0.03]
            },
            'memo': {
                'formats': ['docx', 'doc', 'txt', 'rtf', 'html', 'odt'],
                'weights': [0.35, 0.25, 0.15, 0.1, 0.1, 0.05]
            },
            'contract': {
                'formats': ['pdf', 'docx', 'doc', 'rtf', 'txt', 'html'],
                'weights': [0.35, 0.25, 0.15, 0.1, 0.1, 0.05]
            },
            'legal': {
                'formats': ['pdf', 'docx', 'txt', 'rtf', 'html', 'doc'],
                'weights': [0.5, 0.2, 0.15, 0.08, 0.05, 0.02]
            },
            'report': {
                'formats': ['xlsx', 'docx', 'pdf', 'pptx', 'html', 'txt'],
                'weights': [0.3, 0.25, 0.2, 0.15, 0.07, 0.03]
            },
            'other': {
                'formats': ['txt', 'html', 'docx', 'pdf', 'rtf', 'odt'],
                'weights': [0.25, 0.2, 0.2, 0.15, 0.1, 0.1]
            }
        }
        
        print("\n📋 Generating documents by category:")
        
        for category, count in distribution.items():
            print(f"  📁 {category.capitalize()}: {count} documents")
            
            formats = category_formats[category]['formats']
            weights = category_formats[category]['weights']
            
            for i in range(count):
                fmt = random.choices(formats, weights=weights)[0]
                filename = f"{category}_{i+1:03d}.{fmt}"
                
                try:
                    # Generate file using existing methods for now
                    if fmt == 'pdf':
                        if hasattr(self, f'generate_{category}_pdf'):
                            getattr(self, f'generate_{category}_pdf')(filename)
                        else:
                            # Fallback method
                            self.generate_simple_pdf(filename, category.upper(), f"{category.title()} document")
                    elif fmt == 'docx':
                        if hasattr(self, f'generate_{category}_docx'):
                            getattr(self, f'generate_{category}_docx')(filename)
                        else:
                            self.generate_simple_docx(filename, category.upper(), f"{category.title()} document")
                    elif fmt == 'txt':
                        if hasattr(self, f'generate_{category}_txt'):
                            getattr(self, f'generate_{category}_txt')(filename)
                        else:
                            self.generate_simple_txt(filename, category.upper(), f"{category.title()} document")
                    elif fmt == 'html':
                        self.generate_simple_html(filename, category.upper(), f"{category.title()} document")
                    elif fmt == 'rtf':
                        self.generate_simple_rtf(filename, category.upper(), f"{category.title()} document")
                    elif fmt == 'doc':
                        self.generate_simple_doc(filename, category.upper(), f"{category.title()} document")
                    elif fmt == 'odt':
                        if self.generate_simple_odt(filename, category.upper(), f"{category.title()} document") is None:
                            # Fallback to RTF
                            fmt = 'rtf'
                            filename = filename.replace('.odt', '.rtf')
                            self.generate_simple_rtf(filename, category.upper(), f"{category.title()} document")
                    elif fmt == 'pptx':
                        if self.generate_simple_pptx(filename, category.upper(), f"{category.title()} presentation") is None:
                            # Fallback to PDF
                            fmt = 'pdf'
                            filename = filename.replace('.pptx', '.pdf')
                            self.generate_simple_pdf(filename, category.upper(), f"{category.title()} document")
                    elif fmt == 'xlsx':
                        if self.generate_simple_xlsx(filename, category.upper(), f"{category.title()} data") is None:
                            # Fallback to TXT
                            fmt = 'txt'
                            filename = filename.replace('.xlsx', '.txt')
                            self.generate_simple_txt(filename, category.upper(), f"{category.title()} document")
                    elif fmt == 'xls':
                        if self.generate_simple_xls(filename, category.upper(), f"{category.title()} data") is None:
                            # Fallback to TXT
                            fmt = 'txt'
                            filename = filename.replace('.xls', '.txt')
                            self.generate_simple_txt(filename, category.upper(), f"{category.title()} document")
                    
                    format_count[fmt] += 1
                    generated_files.append(filename)
        
                except Exception as e:
                    print(f"⚠️  Error generating {filename}: {e}")
                    # Generate fallback TXT file
                    txt_filename = f"{category}_{i+1:03d}.txt"
                    try:
                        self.generate_simple_txt(txt_filename, category.upper(), f"{category.title()} document")
                        format_count['txt'] += 1
                        generated_files.append(txt_filename)
                    except:
                        print(f"❌ Failed to generate fallback for {category}_{i+1:03d}")
        
        # Print comprehensive summary
        print(f"\n✅ Successfully generated {len(generated_files)} diverse documents!")
        print(f"\n📊 FORMAT DISTRIBUTION:")
        for fmt, count in format_count.items():
            if count > 0:
                icon = {'pdf': '📄', 'docx': '📝', 'txt': '📋', 'doc': '📄', 'html': '🌐', 
                       'rtf': '📄', 'odt': '📄', 'pptx': '📊', 'ppt': '📊', 'xlsx': '📊', 'xls': '📊'}
                print(f"   {icon.get(fmt, '📄')} {fmt.upper()}: {count} files")
        
        print(f"\n📂 Documents organized by format in '{self.output_dir}/' folder")
        print(f"🎯 Ready for classification testing with enhanced diversity!")
        
        return format_count

    # ===========================================
    # SIMPLE GENERATORS FOR NEW FORMATS
    # ===========================================

    def generate_simple_html(self, filename: str, doc_type: str, description: str):
        """Generate a simple HTML document."""
        filepath = os.path.join(self.output_dir, 'html', filename)
        company = random.choice(self.companies)
        date = self.random_date(90).strftime('%B %d, %Y')
        
        content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>{doc_type}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; }}
        .header {{ text-align: center; border-bottom: 2px solid #333; }}
        .content {{ margin: 20px 0; }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{company}</h1>
        <h2>{doc_type}</h2>
        <p>Date: {date}</p>
    </div>
    <div class="content">
        <p>{description}</p>
        <p>This document contains important business information and procedures.</p>
        <p>Generated on {datetime.now().strftime('%Y-%m-%d')}</p>
    </div>
</body>
</html>"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)

    def generate_simple_rtf(self, filename: str, doc_type: str, description: str):
        """Generate a simple RTF document."""
        filepath = os.path.join(self.output_dir, 'rtf', filename)
        company = random.choice(self.companies)
        date = self.random_date(90).strftime('%B %d, %Y')
        
        rtf_content = r"""{{\rtf1\ansi\deff0 {{\fonttbl{{\f0 Arial;}}}}
\f0\fs24 """ + company + r"""
\par """ + doc_type + r"""
\par Date: """ + date + r"""
\par 
\par """ + description + r"""
\par This document contains important business information.
\par Generated: """ + datetime.now().strftime('%Y-%m-%d') + r"""
\par }}"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(rtf_content)

    def generate_simple_doc(self, filename: str, doc_type: str, description: str):
        """Generate a simple DOC document (RTF format for compatibility)."""
        filepath = os.path.join(self.output_dir, 'doc', filename)
        company = random.choice(self.companies)
        date = self.random_date(90).strftime('%B %d, %Y')
        
        rtf_content = r"""{{\rtf1\ansi\deff0 {{\fonttbl{{\f0 Arial;}}}}
\f0\fs24 """ + company + r"""
\par """ + doc_type + r"""
\par Date: """ + date + r"""
\par 
\par """ + description + r"""
\par This document contains important business information.
\par Format: Legacy DOC Compatible
\par Generated: """ + datetime.now().strftime('%Y-%m-%d') + r"""
\par }}"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(rtf_content)

    def generate_simple_odt(self, filename: str, doc_type: str, description: str):
        """Generate a simple ODT document."""
        if OpenDocumentText is None:
            return None
            
        filepath = os.path.join(self.output_dir, 'odt', filename)
        company = random.choice(self.companies)
        date = self.random_date(90).strftime('%B %d, %Y')
        
        try:
            doc = OpenDocumentText()
            
            title = H(outlinelevel=1, text=doc_type)
            doc.text.addElement(title)
            
            company_para = P(text=f"Company: {company}")
            doc.text.addElement(company_para)
            
            date_para = P(text=f"Date: {date}")
            doc.text.addElement(date_para)
            
            doc.text.addElement(P())
            
            desc_para = P(text=description)
            doc.text.addElement(desc_para)
            
            info_para = P(text="This document contains important business information.")
            doc.text.addElement(info_para)
            
            doc.save(filepath)
            return filepath
        except Exception as e:
            print(f"Error generating ODT: {e}")
            return None

    def generate_simple_pptx(self, filename: str, doc_type: str, description: str):
        """Generate a simple PPTX presentation."""
        if Presentation is None:
            return None
            
        filepath = os.path.join(self.output_dir, 'pptx', filename)
        company = random.choice(self.companies)
        date = self.random_date(90).strftime('%B %d, %Y')
        
        try:
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = doc_type
            subtitle.text = f"{company}\n{date}"
            
            # Content slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Overview"
            content.text = f"• {description}\n• Contains important business information\n• Generated for classification testing"
            
            prs.save(filepath)
            return filepath
        except Exception as e:
            print(f"Error generating PPTX: {e}")
            return None

    def generate_simple_xlsx(self, filename: str, doc_type: str, description: str):
        """Generate a simple XLSX spreadsheet."""
        if pd is None or openpyxl is None:
            return None
            
        filepath = os.path.join(self.output_dir, 'xlsx', filename)
        company = random.choice(self.companies)
        date = self.random_date(90).strftime('%B %d, %Y')
        
        try:
            data = {
                'Field': ['Document Type', 'Company', 'Date', 'Description', 'Status'],
                'Value': [doc_type, company, date, description, 'Active']
            }
            df = pd.DataFrame(data)
            df.to_excel(filepath, index=False)
            return filepath
        except Exception as e:
            print(f"Error generating XLSX: {e}")
            return None

    def generate_simple_xls(self, filename: str, doc_type: str, description: str):
        """Generate a simple XLS spreadsheet."""
        if pd is None:
            return None
            
        filepath = os.path.join(self.output_dir, 'xls', filename)
        company = random.choice(self.companies)
        date = self.random_date(90).strftime('%B %d, %Y')
        
        try:
            data = {
                'Field': ['Type', 'Company', 'Date'],
                'Value': [doc_type, company, date]
            }
            df = pd.DataFrame(data)
            df.to_excel(filepath, index=False, engine='xlwt')
            return filepath
        except Exception as e:
            print(f"Error generating XLS: {e}")
            return None

    def generate_simple_pdf(self, filename: str, doc_type: str, description: str):
        """Generate a simple PDF document."""
        filepath = os.path.join(self.output_dir, 'pdf', filename)
        doc = SimpleDocTemplate(filepath, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        company = random.choice(self.companies)
        date = self.random_date(90).strftime('%B %d, %Y')
        
        story.append(Paragraph(company, styles['Title']))
        story.append(Paragraph(doc_type, styles['Heading1']))
        story.append(Spacer(1, 20))
        story.append(Paragraph(f"Date: {date}", styles['Normal']))
        story.append(Spacer(1, 20))
        story.append(Paragraph(description, styles['Normal']))
        story.append(Paragraph("This document contains important business information.", styles['Normal']))
        
        doc.build(story)

    def generate_simple_docx(self, filename: str, doc_type: str, description: str):
        """Generate a simple DOCX document."""
        filepath = os.path.join(self.output_dir, 'docx', filename)
        doc = DocxDocument()
        
        company = random.choice(self.companies)
        date = self.random_date(90).strftime('%B %d, %Y')
        
        doc.add_heading(doc_type, 0)
        doc.add_paragraph(f"Company: {company}")
        doc.add_paragraph(f"Date: {date}")
        doc.add_paragraph("")
        doc.add_paragraph(description)
        doc.add_paragraph("This document contains important business information.")
        
        doc.save(filepath)

    def generate_simple_txt(self, filename: str, doc_type: str, description: str):
        """Generate a simple TXT document."""
        filepath = os.path.join(self.output_dir, 'txt', filename)
        company = random.choice(self.companies)
        date = self.random_date(90).strftime('%B %d, %Y')
        
        content = f"""{doc_type}

Company: {company}
Date: {date}

{description}

This document contains important business information and procedures.
Generated on {datetime.now().strftime('%Y-%m-%d')}
"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)

def main():
    """Main function to generate diverse sample documents."""
    print("🌟 Welcome to Diverse Document Generator!")
    print("This will create realistic PDF, TXT, and DOCX files for testing.\n")
    
    try:
        num_docs = int(input("Enter number of documents to generate (default: 60): ") or "60")
    except ValueError:
        num_docs = 60
    
    generator = DiverseDocumentGenerator()
    format_counts = generator.generate_all_diverse_documents(num_docs)
    
    print(f"\n🎉 Document generation complete!")
    print(f"📂 All documents saved to 'diverse_sample_documents/' folder")
    print(f"\n💡 To test with these documents:")
    print(f"   python main.py diverse_sample_documents")
    print(f"   python main.py --dashboard")
    print(f"   python launch_dashboard.py")

if __name__ == "__main__":
    main() 